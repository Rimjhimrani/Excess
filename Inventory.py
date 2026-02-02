import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime
import logging
import pickle
import base64
import uuid
import io
import os
import re
import smtplib
import random
from email.mime.text import MIMEText
from typing import Union, Any, Optional, List, Dict
from decimal import Decimal, InvalidOperation
from collections import Counter
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Set page configuration
st.set_page_config(
    page_title="Inventory Management System",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)


# Custom CSS for better styling
st.markdown("""
<style>
.graph-description {
    background-color: #f0f2f6;
    padding: 10px;
    border-radius: 5px;
    margin-bottom: 20px;
    font-style: italic;
    border-left: 4px solid #1f77b4;
}

.metric-container {
    background-color: #ffffff;
    padding: 20px;
    border-radius: 10px;
    box-shadow: 0 2px 4px rgba(0,0,0,0.1);
}

.status-card {
    padding: 15px;
    border-radius: 8px;
    margin: 10px 0;
}

.status-excess {
    background-color: #ffebee;
    border-left: 4px solid #f44336;
}

.status-short {
    background-color: #fff3e0;
    border-left: 4px solid #ff9800;
}

.status-normal {
    background-color: #e8f5e8;
    border-left: 4px solid #4caf50;
}

.success-box {
    background-color: #d4edda;
    border: 1px solid #c3e6cb;
    border-radius: 5px;
    padding: 15px;
    margin: 10px 0;
}

.lock-button {
    background-color: #28a745;
    color: white;
    padding: 10px 20px;
    border-radius: 5px;
    border: none;
    font-weight: bold;
}

.switch-user-button {
    background-color: #007bff;
    color: white;
    padding: 8px 16px;
    border-radius: 5px;
    border: none;
    font-weight: bold;
    margin: 5px 0;
}
</style>
""", unsafe_allow_html=True)

class DataPersistence:
    """SaaS Persistence: Separates files by Company ID and handles session state"""
    
    @staticmethod
    def get_path(company_id, filename):
        if not os.path.exists('data'): 
            os.makedirs('data')
        return f"data/{company_id}_{filename}"

    # --- DISK METHODS (Saves to server) ---
    @staticmethod
    def save_settings(company_id, tolerance, ideal_days):
        path = DataPersistence.get_path(company_id, "settings.pkl")
        with open(path, 'wb') as f:
            pickle.dump({'admin_tolerance': tolerance, 'ideal_inventory_days': ideal_days}, f)
            
    @staticmethod
    def load_settings(company_id):
        path = DataPersistence.get_path(company_id, "settings.pkl")
        if os.path.exists(path):
            with open(path, 'rb') as f:
                return pickle.load(f)
        return None

    @staticmethod
    def save_to_disk(company_id, data, locked=True):
        # This ensures Company A's data goes into "data/COMPANYA_pfep_master.pkl"
        path = DataPersistence.get_path(company_id, "pfep_master.pkl")
        save_obj = {'payload': data, 'timestamp': datetime.now(), 'is_locked': locked}
        with open(path, 'wb') as f:
            pickle.dump(save_obj, f)
            
    @staticmethod
    def load_from_disk(company_id):
        path = DataPersistence.get_path(company_id, "pfep_master.pkl")
        if os.path.exists(path):
            with open(path, 'rb') as f:
                obj = pickle.load(f)
                if isinstance(obj, dict) and 'payload' in obj:
                    return obj['payload'], obj.get('is_locked', False), obj.get('timestamp')
                return obj, False, None
        return None, False, None

    @staticmethod
    def delete_from_disk(company_id):
        path = DataPersistence.get_path(company_id, "pfep_master.pkl")
        if os.path.exists(path): 
            os.remove(path)

    # --- SESSION METHODS (Fixes your AttributeError) ---
    @staticmethod
    def save_data_to_session_state(key, data):
        """Saves data to session state with a timestamp"""
        st.session_state[key] = {
            'data': data,
            'timestamp': datetime.now(),
            'saved': True
        }
    
    @staticmethod
    def load_data_from_session_state(key):
        """Loads actual data from the session state container"""
        container = st.session_state.get(key)
        if container and isinstance(container, dict):
            return container.get('data')
        return container # Return raw if not a dict
    
    @staticmethod
    def get_data_timestamp(key):
        """Retrieves the timestamp for PPT generation"""
        container = st.session_state.get(key)
        if container and isinstance(container, dict):
            return container.get('timestamp')
        return None
        
class InventoryAnalyzer:
    """Enhanced inventory analysis with comprehensive reporting"""
    
    def __init__(self):
        self.debug = False
        self.persistence = DataPersistence()
        self.status_colors = {
            'Within Norms': '#4CAF50',    # Green
            'Excess Inventory': '#2196F3', # Blue
            'Short Inventory': '#F44336'   # Red
        }
        
    def analyze_inventory(self, pfep_data, current_inventory, tolerance=None):
        results = []
        # Normalize and create lookup dictionaries (STRICT STRIP AND UPPER)
        pfep_dict = {str(item['Part_No']).strip().upper(): item for item in pfep_data}
        inventory_dict = {str(item['Part_No']).strip().upper(): item for item in current_inventory}
        
        # Debug info for console
        logger.info(f"Analyzer: PFEP has {len(pfep_dict)} keys. Inventory has {len(inventory_dict)} keys.")

        for part_no, inventory_item in inventory_dict.items():
            pfep_item = pfep_dict.get(part_no)
            
            if not pfep_item:
                continue  # Skip parts not in PFEP master
                
            try:
                current_qty = float(inventory_item.get('Current_QTY', 0))
                unit_price = float(pfep_item.get('unit_price', 100))
                avg_cons = float(pfep_item.get('AVG CONSUMPTION/DAY', 0))
                rm_qty = float(pfep_item.get('RM_IN_QTY', 0))
                
                # Calculation Logic
                upper_bound = np.ceil(rm_qty * (1 + tolerance / 100))
                lower_bound = np.ceil(rm_qty * (1 - tolerance / 100))
                
                deviation_qty = current_qty - upper_bound
                deviation_value = deviation_qty * unit_price
                
                if current_qty < lower_bound:
                    status = 'Short Inventory'
                elif current_qty > upper_bound:
                    status = 'Excess Inventory'
                else:
                    status = 'Within Norms'

                results.append({
                    'PART NO': part_no,
                    'PART DESCRIPTION': pfep_item.get('Description', 'N/A'),
                    'Vendor Name': pfep_item.get('Vendor_Name', 'Unknown'),
                    'UNIT PRICE': unit_price,
                    'AVG CONSUMPTION/DAY': avg_cons,
                    'RM Norm - In Qty': rm_qty,
                    'Revised Norm Qty': upper_bound,
                    'Current Inventory - Qty': current_qty,
                    'Current Inventory - VALUE': current_qty * unit_price,
                    'Stock Deviation Qty w.r.t Revised Norm': deviation_qty,
                    'Stock Deviation Value': deviation_value,
                    'Status': status,
                    'INVENTORY REMARK STATUS': status
                })
            except Exception as e:
                logger.error(f"Error analyzing part {part_no}: {e}")
                
        return results
        
    def get_vendor_summary(self, processed_data):
        """Summarize inventory by vendor using actual Stock_Value field from the file."""
        from collections import defaultdict
        summary = defaultdict(lambda: {
            'total_parts': 0,
            'short_parts': 0,
            'excess_parts': 0,
            'normal_parts': 0,
            'total_value': 0.0,
            'excess_value_above_norm': 0.0,  # New field for excess value above norm
            'short_value_below_norm': 0.0    # New field for short value below norm
        })
        for item in processed_data:
            vendor = item.get('Vendor Name', 'Unknown')
            status = item.get('INVENTORY REMARK STATUS', 'Unknown')
            stock_value = item.get('Stock_Value') or item.get('Current Inventory - VALUE') or 0
            # Get current quantity and norm values for calculating excess/short amounts
            current_qty = item.get('Current Inventory - QTY', 0)
            norm_qty = item.get('Norm', 0)
            unit_price = 0
            try:
                stock_value = float(stock_value)
                current_qty = float(current_qty) if current_qty else 0
                norm_qty = float(norm_qty) if norm_qty else 0
                # Calculate unit price from stock value and current quantity
                if current_qty > 0:
                    unit_price = stock_value / current_qty
            except (ValueError, TypeError):
                stock_value = 0.0
                current_qty = 0.0
                norm_qty = 0.0
                unit_price = 0.0
            summary[vendor]['total_parts'] += 1
            summary[vendor]['total_value'] += stock_value
            if status == "Short Inventory":
                summary[vendor]['short_parts'] += 1
                # Calculate value of shortage (norm - current) * unit_price
                if norm_qty > current_qty and unit_price > 0:
                    short_value = (norm_qty - current_qty) * unit_price
                    summary[vendor]['short_value_below_norm'] += short_value
            elif status == "Excess Inventory":
                summary[vendor]['excess_parts'] += 1
                # Calculate value of excess (current - norm) * unit_price
                if current_qty > norm_qty and unit_price > 0:
                    excess_value = (current_qty - norm_qty) * unit_price
                    summary[vendor]['excess_value_above_norm'] += excess_value
            elif status == "Within Norms":
                summary[vendor]['normal_parts'] += 1
        return summary
        
    def show_vendor_chart_by_status(self, processed_data, status_filter, chart_title, chart_key, color, value_format='lakhs', top_n=10):
        """Show top N vendors by deviation value and count of excess/short parts."""
        # Filter by inventory status
        filtered = [item for item in processed_data if item.get('INVENTORY REMARK STATUS') == status_filter]
        vendor_totals = {}
        vendor_counts = {}
        for item in filtered:
            vendor = item.get('Vendor Name', 'Unknown')
            try:
                current_qty = float(item.get('Current Inventory - Qty', 0) or 0)
                norm_qty = float(item.get('Revised Norm Qty', 0) or 0)
                stock_value = float(item.get('Current Inventory - VALUE', 0) or 0)
                unit_price = float(item.get('UNIT PRICE', 0) or 0)
                if unit_price == 0 and current_qty > 0:
                    unit_price = stock_value / current_qty
                if status_filter == "Excess Inventory" and current_qty > norm_qty:
                    deviation_value = (current_qty - norm_qty) * unit_price
                    vendor_totals[vendor] = vendor_totals.get(vendor, 0.0) + deviation_value
                    vendor_counts[vendor] = vendor_counts.get(vendor, 0) + 1
                elif status_filter == "Short Inventory" and norm_qty > current_qty:
                    deviation_value = (norm_qty - current_qty) * unit_price
                    vendor_totals[vendor] = vendor_totals.get(vendor, 0.0) + deviation_value
                    vendor_counts[vendor] = vendor_counts.get(vendor, 0) + 1
            except (ValueError, TypeError, ZeroDivisionError):
                continue
        if not vendor_totals:
            st.info(f"No vendors found in '{status_filter}' status.")
            return
        # Combine totals and counts
        combined = [
            (vendor, vendor_totals[vendor], vendor_counts.get(vendor, 0))
            for vendor in vendor_totals
        ]
        # Sort and take top N by value (Updated Logic Here)
        top_vendors = sorted(combined, key=lambda x: x[1], reverse=True)[:top_n]
        vendor_names = [v[0] for v in top_vendors]
        raw_values = [v[1] for v in top_vendors]
        counts = [v[2] for v in top_vendors]
        
        # Update Title to reflect dynamic number
        dynamic_title = chart_title.replace("Top 10", f"Top {top_n}")

        # Value formatting
        if value_format == 'lakhs':
            values = [v / 100000 for v in raw_values]
            y_axis_title = "Value (₹ Lakhs)"
            hover_suffix = "L"
            tick_suffix = "L"
        elif value_format == 'millions':
            values = [v / 1000000 for v in raw_values]
            y_axis_title = "Value (Millions)"
            hover_suffix = "M"
            tick_suffix = "M"
        elif value_format == 'crores':
            values = [v / 10000000 for v in raw_values]
            y_axis_title = "Value (₹ Crores)"
            hover_suffix = "Cr"
            tick_suffix = "Cr"
        else:
            values = raw_values
            y_axis_title = "Value (₹)"
            hover_suffix = ""
            tick_suffix = ""
            
        if status_filter == "Excess Inventory":
            y_axis_title = y_axis_title.replace("Value", "Excess Value Above Norm")
            hover_label = "Excess Value"
        elif status_filter == "Short Inventory":
            y_axis_title = y_axis_title.replace("Value", "Short Value Below Norm")
            hover_label = "Short Value"
        else:
            hover_label = "Stock Value"
        # Create chart
        fig = go.Figure()
        fig.add_trace(go.Bar(
            x=vendor_names,
            y=values,
            marker_color=color,
            text=[f"{hover_label}: ₹{v:,.1f}{hover_suffix}<br>Parts: {c}" for v, c in zip(values, counts)],
            hovertemplate='<b>%{x}</b><br>%{text}<extra></extra>',
        ))
        fig.update_layout(
            title=dynamic_title,
            xaxis_title="Vendor",
            yaxis_title=y_axis_title,
            showlegend=False,
            yaxis=dict(tickformat=".1f", ticksuffix=tick_suffix),
        )
        st.plotly_chart(fig, use_container_width=True, key=chart_key)

    def safe_float_convert(self, value, default=0.0):
        try:
            return float(value)
        except (ValueError, TypeError):
            return default

class InventoryManagementSystem:
    """Main application class"""
    
    def __init__(self):
        self.debug = False
        self.analyzer = InventoryAnalyzer()
        self.persistence = DataPersistence()
        self.initialize_session_state()
        self.status_colors = {
            'Within Norms': '#4CAF50',    # Green
            'Excess Inventory': '#2196F3', # Blue
            'Short Inventory': '#F44336'   # Red
        }
        
    def initialize_session_state(self):
        """Initializes keys and loads data ONLY when the company changes"""
        # 1. Basic initialization
        if 'user_role' not in st.session_state: st.session_state.user_role = None
        if 'company_id' not in st.session_state: st.session_state.company_id = None
        if 'last_loaded_company' not in st.session_state: st.session_state.last_loaded_company = None
    
        comp_id = st.session_state.get('company_id')
        if not comp_id:
            return 

        # 2. CHECK: Has the company changed since the last run?
        # If comp_id is the same as last_loaded_company, STOP here so we don't wipe out "Run Analysis" results
        if comp_id == st.session_state.last_loaded_company:
             return

        # 3. IF COMPANY IS DIFFERENT (New Login):
        # Clear old company data and load new company data
        st.session_state['persistent_inventory_data'] = None
        st.session_state['persistent_analysis_results'] = None
    
        # Load settings for the new company
        saved_settings = DataPersistence.load_settings(comp_id)
        st.session_state.admin_tolerance = saved_settings['admin_tolerance'] if saved_settings else 30
        ideal = saved_settings['ideal_inventory_days'] if saved_settings else 30
        st.session_state.user_preferences = {'ideal_inventory_days': ideal, 'chart_theme': 'plotly'}

        # Load master data from disk for the new company
        disk_data, is_locked, disk_ts = DataPersistence.load_from_disk(comp_id)
        if disk_data:
            st.session_state['persistent_pfep_data'] = {
                'data': disk_data, 
                'timestamp': disk_ts
            }
            st.session_state['persistent_pfep_locked'] = is_locked
        else:
            st.session_state['persistent_pfep_data'] = None
            st.session_state['persistent_pfep_locked'] = False

        # Mark this company as loaded so we don't repeat this on next rerun
        st.session_state.last_loaded_company = comp_id
    
    def safe_print(self, message):
        """Safely print to streamlit or console"""
        try:
            st.write(message)
        except NameError:
            print(message)
    
    def safe_error(self, message):
        """Safely show error in streamlit or console"""
        try:
            st.error(message)
        except NameError:
            print(f"ERROR: {message}")
    
    def safe_warning(self, message):
        """Safely show warning in streamlit or console"""
        try:
            st.warning(message)
        except NameError:
            print(f"WARNING: {message}")
    
    def safe_float_convert(self, value: Any) -> float:
        """Safely convert value to float"""
        if value is None or pd.isna(value) or value == '':
            return 0.0
        
        if isinstance(value, (int, float)):
            return float(value)
        
        try:
            # Clean string values
            if isinstance(value, str):
                value = value.strip()
                # Remove currency symbols and commas
                value = re.sub(r'[₹$€£,]', '', value)
                # Handle percentage
                if '%' in value:
                    return float(value.replace('%', '')) / 100
            return float(value)
        except (ValueError, TypeError) as e:
            # Debug problematic values
            if hasattr(self, 'debug') and self.debug:
                st.warning(f"⚠️ Could not convert '{value}' to float: {e}")
            return 0.0
            
    def safe_int_convert(self, value: Any) -> int:
        """Convert value to integer using safe float conversion."""
        float_result = self.safe_float_convert(value)
        return int(float_result)
            
    def create_top_parts_chart(self, data, status_filter, bar_color, key):
        """Top 10 parts chart by status — shows EXCESS/SHORT VALUE vs NORM VALUE comparison"""
        df = pd.DataFrame(data)
        # ✅ Check required columns
        if 'PART NO' not in df.columns or 'Stock Deviation Value' not in df.columns:
            st.warning("⚠️ Required columns missing for top parts chart.")
            return
        # ✅ Filter by status
        df = df[df['INVENTORY REMARK STATUS'] == status_filter]
    
        if status_filter == "Excess Inventory":
            # For excess inventory, show only positive deviation values (excess amount)
            df = df[df['Stock Deviation Value'] > 0]
            df = df.sort_values(by='Stock Deviation Value', ascending=False).head(10)
            chart_title = "Top 10 Excess Inventory Parts - Norm vs Actual (₹ in Lakhs)"
            y_title = "Inventory Value (₹ Lakhs)"
        elif status_filter == "Short Inventory":
            # For short inventory, show absolute value of negative deviation
            df = df[df['Stock Deviation Value'] < 0]
            df['Abs_Deviation_Value'] = abs(df['Stock Deviation Value'])
            df = df.sort_values(by='Abs_Deviation_Value', ascending=False).head(10)
            chart_title = "Top 10 Short Inventory Parts - Norm vs Actual (₹ in Lakhs)"
            y_title = "Inventory Value (₹ Lakhs)"
        else:
            st.info(f"No chart available for '{status_filter}' status.")
            return
        if df.empty:
            st.info(f"No data found for '{status_filter}' parts.")
            return
        # ✅ Prepare chart data with norm vs actual comparison
        if status_filter == "Excess Inventory":
            # Calculate norm value (current stock value - excess)
            df['Current_Stock_Value'] = df['CURRENT STOCK VALUE'] if 'CURRENT STOCK VALUE' in df.columns else 0
            df['Norm_Value'] = df['Current_Stock_Value'] - df['Stock Deviation Value']
            df['Excess_Value'] = df['Stock Deviation Value']
            # Convert to lakhs
            df['Norm_Value_Lakh'] = df['Norm_Value'] / 100000
            df['Current_Value_Lakh'] = df['Current_Stock_Value'] / 100000
            df['Deviation_Value_Lakh'] = df['Excess_Value'] / 100000
        else:  # Short Inventory
            # Calculate what the stock value should be (current + shortage)
            df['Current_Stock_Value'] = df['CURRENT STOCK VALUE'] if 'CURRENT STOCK VALUE' in df.columns else 0
            df['Norm_Value'] = df['Current_Stock_Value'] + df['Abs_Deviation_Value']
            df['Shortage_Value'] = df['Abs_Deviation_Value']
        
            # Convert to lakhs
            df['Norm_Value_Lakh'] = df['Norm_Value'] / 100000
            df['Current_Value_Lakh'] = df['Current_Stock_Value'] / 100000
            df['Deviation_Value_Lakh'] = df['Shortage_Value'] / 100000
        # Prepare labels and hover text
        df['PART_DESC_NO'] = df['PART DESCRIPTION'].astype(str) + " (" + df['PART NO'].astype(str) + ")"
    
        if status_filter == "Excess Inventory":
            df['HOVER_TEXT_NORM'] = df.apply(lambda row: (
                f"Description: {row.get('PART DESCRIPTION', 'N/A')}<br>"
                f"Part No: {row.get('PART NO')}<br>"
                f"Norm Value: ₹{row['Norm_Value']:,.0f}<br>"
                f"Type: Within Norm Limit"
            ), axis=1)
        
            df['HOVER_TEXT_CURRENT'] = df.apply(lambda row: (
                f"Description: {row.get('PART DESCRIPTION', 'N/A')}<br>"
                f"Part No: {row.get('PART NO')}<br>"
                f"Current Value: ₹{row['Current_Stock_Value']:,.0f}<br>"
                f"Excess Value: ₹{row['Excess_Value']:,.0f}<br>"
                f"Type: Current Stock (Excess)"
            ), axis=1)
        else:  # Short Inventory
            df['HOVER_TEXT_NORM'] = df.apply(lambda row: (
                f"Description: {row.get('PART DESCRIPTION', 'N/A')}<br>"
                f"Part No: {row.get('PART NO')}<br>"
                f"Required Value: ₹{row['Norm_Value']:,.0f}<br>"
                f"Type: Required Norm Level"
            ), axis=1)
        
            df['HOVER_TEXT_CURRENT'] = df.apply(lambda row: (
                f"Description: {row.get('PART DESCRIPTION', 'N/A')}<br>"
                f"Part No: {row.get('PART NO')}<br>"
                f"Current Value: ₹{row['Current_Stock_Value']:,.0f}<br>"
                f"Shortage Value: ₹{row['Shortage_Value']:,.0f}<br>"
                f"Type: Current Stock (Short)"
            ), axis=1)
        fig = go.Figure()
        # Add norm value bars
        fig.add_trace(go.Bar(
            x=df['PART_DESC_NO'],
            y=df['Norm_Value_Lakh'],
            name='Norm Value' if status_filter == "Excess Inventory" else 'Required Value',
            marker_color='#4CAF50',  # Green for norm
            hovertemplate='<b>%{x}</b><br>%{customdata}<extra></extra>',
            customdata=df['HOVER_TEXT_NORM']
        ))
        # Add current value bars
        fig.add_trace(go.Bar(
            x=df['PART_DESC_NO'],
            y=df['Current_Value_Lakh'],
            name='Current Value',
            marker_color=bar_color,
            hovertemplate='<b>%{x}</b><br>%{customdata}<extra></extra>',
            customdata=df['HOVER_TEXT_CURRENT']
        ))
    
        # ✅ Update chart formatting
        fig.update_layout(
            title=chart_title,
            xaxis_title="Part Description (Part No)",
            yaxis_title=y_title,
            xaxis_tickangle=-45,
            barmode='group',  # Side-by-side bars
            yaxis=dict(
                tickformat=',.1f',
                ticksuffix='L'
            ),
            legend=dict(
                orientation="h",
                yanchor="bottom",
                y=1.02,
                xanchor="right",
                x=1
            ),
            height=600,
            showlegend=True
        )
        st.plotly_chart(fig, use_container_width=True, key=key)
    
        # ✅ Add summary table showing the comparison
        st.subheader(f"📊 Summary - {status_filter}")
        summary_df = df[['PART NO', 'PART DESCRIPTION', 'Current_Value_Lakh', 'Norm_Value_Lakh', 'Deviation_Value_Lakh']].copy()
        summary_df.columns = ['Part No', 'Part Description', 'Current Value (₹L)', 'Norm Value (₹L)', 
                              'Excess Value (₹L)' if status_filter == "Excess Inventory" else 'Shortage Value (₹L)']
        summary_df = summary_df.round(2)
        st.dataframe(summary_df, use_container_width=True)

    def authenticate_user(self):
        st.sidebar.markdown("### 🔐 Corporate Login")
        path = "data/company_registry.pkl"
        if os.path.exists(path):
            with open(path, "rb") as f: registry = pickle.load(f)
        else: registry = {}

        # Handle View States
        if st.session_state.get('changing_password'):
            self.handle_password_change(st.session_state.temp_comp_id, registry)
            return
    
        if st.session_state.get('reset_mode'):
            self.handle_forgot_password_view(registry)
            return

        if st.session_state.user_role is None:
            comp_id = st.sidebar.text_input("Company ID").upper().strip()
        
            # --- FORGOT PASSWORD BUTTON ---
            if st.sidebar.button("❓ Forgot Password"):
                if comp_id in registry:
                    dest_email = registry[comp_id].get("email")
                    if dest_email:
                        otp = str(random.randint(100000, 999999))
                        if self.send_otp_email(dest_email, otp):
                            st.session_state.generated_otp = otp
                            st.session_state.reset_target_id = comp_id
                            st.session_state.reset_mode = True
                            st.rerun()
                    else:
                        st.sidebar.error("No recovery email found for this ID.")
                else:
                    st.sidebar.error("Please enter a valid Company ID first.")

            role = st.sidebar.selectbox("Role", ["Select Role", "Admin", "User"])
            if role == "Admin":
                password = st.sidebar.text_input("Password", type="password")
                if st.sidebar.button("🔑 Login"):
                    if comp_id in registry and password == registry[comp_id]["password"]:
                        if registry[comp_id].get("status") == "FIRST_LOGIN":
                            st.session_state.changing_password = True
                            st.session_state.temp_comp_id = comp_id
                            st.rerun()
                    
                        st.session_state.user_role = "Admin"
                        st.session_state.company_id = comp_id
                        self.initialize_session_state()
                        st.rerun()
                    else:
                        st.sidebar.error("❌ Invalid ID or Password")
        
            elif role == "User":
                if st.sidebar.button("👤 Enter as User"):
                    if comp_id in registry:
                        st.session_state.user_role = "User"
                        st.session_state.company_id = comp_id
                        self.initialize_session_state()
                        st.rerun()
                    else:
                        st.sidebar.error("❌ Company ID not recognized")
        else:
            # LOGGED IN VIEW
            st.sidebar.success(f"🏢 **{st.session_state.company_id}**")
            st.sidebar.info(f"👤 Role: {st.session_state.user_role}")
            self.display_data_status()

            if st.session_state.user_role == "Admin":
                pfep_locked = st.session_state.get("persistent_pfep_locked", False)
                if pfep_locked:
                    if st.sidebar.button("🔄 Switch to User View"):
                        st.session_state.user_role = "User"
                        st.rerun()

                with st.sidebar.expander("⚙️ Analysis Settings"):
                    new_tol = st.selectbox("Tolerance %", [0, 10, 20, 30, 40, 50], 
                                           index=[0, 10, 20, 30, 40, 50].index(st.session_state.admin_tolerance))
                    new_ideal = st.number_input("Ideal Days", value=st.session_state.user_preferences['ideal_inventory_days'])
                    if st.button("🔒 Lock Settings"):
                        st.session_state.admin_tolerance = new_tol
                        st.session_state.user_preferences['ideal_inventory_days'] = new_ideal
                        DataPersistence.save_settings(st.session_state.company_id, new_tol, new_ideal)
                        st.success("Settings saved for Company")

            st.sidebar.markdown("---")
            if st.sidebar.button("🚪 Logout"):
                # This completely wipes the browser memory for this session
                for key in list(st.session_state.keys()):
                    del st.session_state[key]
                st.rerun()
    
        self.developer_console()
    
    def developer_console(self):
        with st.sidebar.expander("🛠️ Developer Console"):
            dev_key = st.text_input("Developer Master Key", type="password")
            if dev_key == "AgiloSaaS2026":
                st.subheader("Register New Corporate Client")
                new_c = st.text_input("New Company ID").upper().strip()
                admin_email = st.text_input("Admin Recovery Email") # Added Email field
            
                if st.button("➕ Register Client"):
                    path = "data/company_registry.pkl"
                    if os.path.exists(path):
                        with open(path, "rb") as f: reg = pickle.load(f)
                    else: reg = {}
                
                    reg[new_c] = {
                        "password": "Welcome@123",
                        "status": "FIRST_LOGIN",
                        "email": admin_email  # Save email to registry
                    }
                
                    if not os.path.exists('data'): os.makedirs('data')
                    with open(path, "wb") as f: pickle.dump(reg, f)
                    st.success(f"✅ {new_c} registered. Recovery email: {admin_email}")

    def handle_password_change(self, comp_id, registry):
        """UI for clients to set their private password"""
        st.title("🛡️ Secure Your Account")
        st.warning(f"Required: Set a private password for **{comp_id}**")
        
        new_p = st.text_input("New Private Password", type="password")
        conf_p = st.text_input("Confirm Password", type="password")
    
        if st.button("💾 Save Password & Login"):
            if new_p == conf_p and len(new_p) >= 6:
                # Update values in memory
                registry[comp_id]["password"] = new_p
                registry[comp_id]["status"] = "ACTIVE"
            
                # Write back to the .pkl file
                with open("data/company_registry.pkl", "wb") as f:
                    pickle.dump(registry, f)
            
                st.success("✅ Password set! Redirecting to dashboard...")
            
                # Log them in automatically
                st.session_state.changing_password = False
                st.session_state.user_role = "Admin"
                st.session_state.company_id = comp_id
                self.initialize_session_state()
                st.rerun()
            else:
                st.error("Passwords must match and be at least 6 characters.")

    def send_otp_email(self, target_email, otp_code):
        # ^^^ YOU MUST HAVE 'self' AS THE FIRST WORD INSIDE THE BRACKETS ^^^
        
        # SETTINGS: Update these with your own SMTP details
        sender_email = "agilomatrix1@gmail.com" 
        sender_password = "kgqg dfiq gelu vqli" # Use a Google App Password
    
        msg = MIMEText(f"Your Inventory System password reset code is: {otp_code}")
        msg['Subject'] = 'Password Reset OTP'
        msg['From'] = sender_email
        msg['To'] = target_email

        try:
            with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
                server.login(sender_email, sender_password)
                server.sendmail(sender_email, target_email, msg.as_string())
            return True
        except Exception as e:
            st.error(f"Mail Error: {e}")
            return False
    def handle_forgot_password_view(self, registry):
        st.title("🔑 Reset Password")
        comp_id = st.session_state.get('reset_target_id')
        if not comp_id:
            st.error("Session expired. Please start over.")
            if st.button("Back to Login"):
                st.session_state.reset_mode = False
                st.rerun()
            return

        st.info(f"An OTP has been sent to the registered email for {comp_id}")
    
        input_otp = st.text_input("Enter 6-Digit OTP")
        new_p = st.text_input("New Password", type="password")
        conf_p = st.text_input("Confirm Password", type="password")
    
        if st.button("💾 Reset Password"):
            if input_otp == st.session_state.get('generated_otp'):
                if new_p == conf_p and len(new_p) >= 6:
                    registry[comp_id]["password"] = new_p
                    with open("data/company_registry.pkl", "wb") as f:
                        pickle.dump(registry, f)
                    st.success("✅ Password updated! Please login.")
                    st.session_state.reset_mode = False
                    st.session_state.generated_otp = None # Clear OTP
                    st.rerun()
                else:
                    st.error("Passwords must match and be 6+ characters.")
            else:
                st.error("Invalid OTP code.")
    
        if st.button("⬅️ Back to Login"):
            st.session_state.reset_mode = False
            st.rerun()
    
    def display_data_status(self):
        """Display current data loading status in sidebar"""
        st.sidebar.markdown("---")
        st.sidebar.markdown("### 📊 Data Status")
        
        # Check persistent PFEP data
        pfep_data = self.persistence.load_data_from_session_state('persistent_pfep_data')
        pfep_locked = st.session_state.get('persistent_pfep_locked', False)
        
        if pfep_data:
            pfep_count = len(pfep_data)
            lock_icon = "🔒" if pfep_locked else "🔓"
            st.sidebar.success(f"✅ PFEP Data: {pfep_count} parts {lock_icon}")
            timestamp = self.persistence.get_data_timestamp('persistent_pfep_data')
            if timestamp:
                st.sidebar.caption(f"Loaded: {timestamp.strftime('%Y-%m-%d %H:%M')}")
        else:
            st.sidebar.error("❌ PFEP Data: Not loaded")
        
        # Check persistent inventory data
        inventory_data = self.persistence.load_data_from_session_state('persistent_inventory_data')
        inventory_locked = st.session_state.get('persistent_inventory_locked', False)
        
        if inventory_data:
            inv_count = len(inventory_data)
            lock_icon = "🔒" if inventory_locked else "🔓"
            st.sidebar.success(f"✅ Inventory: {inv_count} parts {lock_icon}")
            timestamp = self.persistence.get_data_timestamp('persistent_inventory_data')
            if timestamp:
                st.sidebar.caption(f"Loaded: {timestamp.strftime('%Y-%m-%d %H:%M')}")
        else:
            st.sidebar.error("❌ Inventory: Not loaded")
        
        # Analysis results status
        analysis_data = self.persistence.load_data_from_session_state('persistent_analysis_results')
        if analysis_data:
            st.sidebar.info(f"📈 Analysis: {len(analysis_data)} parts analyzed")
    
    def load_sample_pfep_data(self):
        pfep_sample = [
           ["AC0303020106", "FLAT ALUMINIUM PROFILE", 4.000, "V001", "Vendor_A", "Mumbai", "Maharashtra", 2.5],
           ["JJ1010101010", "WINDSHIELD WASHER", 25, "V002", "Vendor_B", "Delhi", "Delhi", 1.8]
           # Add more sample data with consumption values...
        ]
        pfep_data = []
        for row in pfep_sample:
            pfep_data.append({
                'Part_No': row[0],
                'Description': row[1],
                'RM_IN_QTY': self.safe_float_convert(row[2]),
                'Vendor_Code': row[3],
                'Vendor_Name': row[4],
                'City': row[5],
                'State': row[6],
                'Unit_Price': 100,            # 🔁 you can customize this per part
                'RM_IN_DAYS': 7,              # 🔁 default or configurable
                'AVG CONSUMPTION/DAY': self.safe_float_convert(row[7]) if len(row) > 7 else ""  # ✅ Added consumption data
            })
        return pfep_data
    
    def load_sample_current_inventory(self):
        """Load sample current inventory data with consistent fields"""
        current_sample = [
            ["AC0303020106", "FLAT ALUMINIUM PROFILE", 5.230, 496],
            # ... rest of your data
            ["JJ1010101010", "WINDSHIELD WASHER", 33, 495]
        ]
        return [{
            'Part_No': row[0],
            'Description': row[1],
            'Current_QTY': self.safe_float_convert(row[2]),
            'Current Inventory - VALUE': self.safe_float_convert(row[3])
        } for row in current_sample]
    
    def standardize_pfep_data(self, df):
        """Enhanced PFEP data standardization with added Unit_Price and RM_IN_DAYS support"""
        if df is None or df.empty:
            return []
        # Debug: Show original column names
        if self.debug:
            st.write("🔍 DEBUG: Original Excel columns:")
            for i, col in enumerate(df.columns):
                st.write(f"  {i}: '{col}' (type: {type(col)})")
        # Comprehensive column mapping with more variations
        column_mappings = {
            'part_no': [
                'part_no', 'part_number', 'material', 'material_code', 'item_code', 
                'code', 'part no', 'partno', 'Part No', 'Part_No', 'PART_NO',
                'Material', 'Material Code', 'Item Code'
            ],
            'description': [
                'description', 'item_description', 'part_description', 'desc', 
                'part description', 'material_description', 'item desc', 'Part Description',
                'Description', 'DESCRIPTION', 'Material Description'
            ],
            'rm_qty': [
                'rm_in_qty', 'rm_qty', 'required_qty', 'norm_qty', 'target_qty', 
                'rm', 'ri_in_qty', 'rm in qty', 'RM_IN_QTY', 'RM_QTY', 'RM IN QTY',
                'Required Qty', 'Norm Qty', 'Target Qty'
            ],
            'rm_days': [
                'rm_in_days', 'rm days', 'inventory days', 'rmindays', 'RM_IN_DAYS',
                'RM IN DAYS', 'RM Days', 'Inventory Days'
            ],
            'unit_price': [
                'unit_price', 'price', 'unit cost', 'unit rate', 'unitprice', 'Unit Price',
                'UNIT_PRICE', 'UNIT PRICE', 'Price', 'PRICE', 'Unit Cost', 'UNIT_COST',
                'unit_cost', 'Unit Rate', 'UNIT_RATE', 'unit_rate', 'rate', 'Rate', 'RATE',
                'cost', 'Cost', 'COST', 'Unit_Price'
            ],
            'vendor_code': [
                'vendor_code', 'vendor_id', 'supplier_code', 'supplier_id', 'vendor id', 
                'Vendor Code', 'vendor code', 'VENDOR_CODE', 'VENDOR CODE', 'Vendor_Code',
                'Supplier Code', 'SUPPLIER_CODE'
            ],
            'vendor_name': [
                'vendor_name', 'vendor', 'supplier_name', 'supplier', 'Vendor Name', 
                'vendor name', 'VENDOR_NAME', 'VENDOR NAME', 'Vendor_Name',
                'Supplier Name', 'SUPPLIER_NAME', 'Supplier'
            ],
            'avg_consumption_per_day': [
                'avg consumption/day', 'average consumption/day', 'avg_per_day',
                'avg daily usage', 'AVG CONSUMPTION/DAY', 'Average Per Day',
                'avg consumption per day', 'daily consumption', 'consumption per day',
                'AVG_CONSUMPTION_PER_DAY', 'AVERAGE_CONSUMPTION_PER_DAY'
            ],
            'city': ['city', 'location', 'place', 'City', 'CITY', 'Location', 'LOCATION'],
            'state': ['state', 'region', 'province', 'State', 'STATE', 'Region', 'REGION']
        }
        # Create case-insensitive column lookup
        df_columns_lookup = {}
        for col in df.columns:
            if col is not None:
                # Clean column name (remove extra spaces, special characters)
                clean_col = str(col).strip()
                df_columns_lookup[clean_col.lower()] = clean_col
        if self.debug:
            st.write("🔍 DEBUG: Cleaned column lookup:")
            for k, v in df_columns_lookup.items():
                st.write(f"  '{k}' -> '{v}'")
        # Map columns using case-insensitive matching
        mapped_columns = {}
        for key, variations in column_mappings.items():
            found = False
            for variation in variations:
                variation_lower = variation.lower().strip()
                if variation_lower in df_columns_lookup:
                    mapped_columns[key] = df_columns_lookup[variation_lower]
                    if self.debug:
                        st.write(f"✅ Mapped {key} -> '{mapped_columns[key]}' (from variation: '{variation}')")
                    found = True
                    break
            if not found and self.debug:
                st.write(f"❌ No mapping found for {key}")
        # Check for required columns
        if 'part_no' not in mapped_columns:
            st.error("❌ Part Number column not found. Please ensure your file has a Part Number column.")
            return []
        if 'rm_qty' not in mapped_columns:
            st.error("❌ RM Quantity column not found. Please ensure your file has an RM/Required Quantity column.")
            return []
        # Warning for missing unit price
        if 'unit_price' not in mapped_columns:
            st.warning("⚠️ Unit Price column not found. Using default value of 100 for all parts.")
            st.info("💡 Expected Unit Price column names: Unit Price, Price, Unit Cost, Rate, etc.")
            
        if 'avg_consumption_per_day' not in mapped_columns:
            st.warning("⚠️ AVG CONSUMPTION/DAY column not found. This will be left empty.")
            st.info("💡 Expected column names: AVG CONSUMPTION/DAY, Average Per Day, Daily Consumption, etc.")
            
        if self.debug:
            st.write("🔍 DEBUG: Final column mappings:")
            for k, v in mapped_columns.items():
                st.write(f"  {k} -> '{v}'")
        # Standardize data
        standardized_data = []
        for idx, row in df.iterrows():
            try:
                # Extract unit price with detailed debugging
                unit_price_value = 100.0  # Default value
                if 'unit_price' in mapped_columns:
                    raw_price = row[mapped_columns['unit_price']]
                    unit_price_value = self.safe_float_convert(raw_price)
                
                    # Extract AVG CONSUMPTION/DAY with proper handling
                    avg_consumption_value = ""  # Default empty string
                    if 'avg_consumption_per_day' in mapped_columns:
                        raw_consumption = row[mapped_columns['avg_consumption_per_day']]
                        # Handle different data types
                        if pd.notna(raw_consumption) and str(raw_consumption).strip() != '':
                            avg_consumption_value = self.safe_float_convert(raw_consumption)
                        
                item = {
                    'Part_No': str(row[mapped_columns['part_no']]).strip(),
                    'Description': str(row.get(mapped_columns.get('description', ''), '')).strip(),
                    'RM_IN_QTY': self.safe_float_convert(row[mapped_columns['rm_qty']]),
                    'RM_IN_DAYS': self.safe_float_convert(row.get(mapped_columns.get('rm_days', ''), 7)),  # Default 7 days
                    'unit_price': unit_price_value,  # Fixed: use lowercase to match analyzer expectation
                    'Vendor_Code': str(row.get(mapped_columns.get('vendor_code', ''), '')).strip(),
                    'Vendor_Name': str(row.get(mapped_columns.get('vendor_name', ''), 'Unknown')).strip(),
                    'AVG CONSUMPTION/DAY': row.get(mapped_columns.get('avg_consumption_per_day', ''), ''),
                    'City': str(row.get(mapped_columns.get('city', ''), '')).strip(),
                    'State': str(row.get(mapped_columns.get('state', ''), '')).strip()
                }
                # Skip rows with empty part numbers
                if not item['Part_No'] or item['Part_No'].lower() in ['nan', 'none', '']:
                    continue
                standardized_data.append(item)
            except Exception as e:
                if self.debug:
                    st.warning(f"⚠️ Error processing row {idx+1}: {e}")
                continue
        # Summary of data processing
        st.success(f"✅ Processed {len(standardized_data)} PFEP records")
        # Show unit price statistics
        if standardized_data:
            prices = [item['unit_price'] for item in standardized_data if item['unit_price'] > 0]
            if prices:
                avg_price = sum(prices) / len(prices)
                st.info(f"💰 Unit Price Summary: {len(prices)} parts with prices, Average: ₹{avg_price:.2f}")
            else:
                st.warning("⚠️ No valid unit prices found in the data")
            # Show AVG CONSUMPTION/DAY statistics
            consumption_values = [
                item.get('AVG CONSUMPTION/DAY')
                for item in standardized_data
                if item.get('AVG CONSUMPTION/DAY') not in (None, '', 'nan')
            ]
            if consumption_values:
                # Convert strings to floats, skip zeros
                numeric = [self.safe_float_convert(val) for val in consumption_values if self.safe_float_convert(val) > 0]
                if numeric:
                    avg_consumption = sum(numeric) / len(numeric)
                    st.info(f"📊 AVG CONSUMPTION/DAY Summary: {len(numeric)} parts with consumption data, Average: {avg_consumption:.2f}")
                else:
                    st.warning("⚠️ All AVG CONSUMPTION/DAY values were zero or invalid.")
            else:
                st.warning("⚠️ No AVG CONSUMPTION/DAY data found in the file")
        return standardized_data
    
    def standardize_current_inventory(self, df):
        """Standardize current inventory data with full column mappings and debugging."""
        if df is None or df.empty:
            return []
        # 🔁 Add all possible column mappings
        column_mappings = {
            'part_no': ['part_no', 'part_number', 'material', 'material_code', 'item_code', 'code', 'part no', 'Part No'],
            'description': ['description', 'item_description', 'part_description', 'desc','Part Description'],
            'current_qty': ['current_qty', 'qty', 'quantity', 'stock_qty', 'available_qty', 'Current_QTY'],
            'stock_value': ['stock_value', 'value', 'total_value', 'inventory_value', 'Stock Value', 'Stock_Value','Current Inventory - VALUE', 'Current Inventory-VALUE'],
            'uom': ['uom', 'unit', 'unit_of_measure'],
            'location': ['location', 'store', 'warehouse', 'site'],
            'vendor_code': ['vendor_code', 'vendor id', 'supplier_code', 'Vendor Code'],
            'batch': ['batch', 'batch_number', 'lot', 'lot_number']
        }
        df_columns_lower = {col.lower().strip(): col for col in df.columns if col is not None}
        mapped_columns = {}
        for key, variations in column_mappings.items():
            for variation in variations:
                if variation.lower() in df_columns_lower:
                    mapped_columns[key] = df_columns_lower[variation.lower()]
                    break
        # Debug: show mappings
        if self.debug:
            st.write("🔍 DEBUG: Column mappings found:")
            for key, col in mapped_columns.items():
                st.write(f"  {key} → {col}")
        if 'part_no' not in mapped_columns or 'current_qty' not in mapped_columns:
            st.error("❌ Required columns not found. Please ensure your file has Part Number and Current Quantity columns.")
            return []
        standardized_data = []
        for i, (_, row) in enumerate(df.iterrows()):
            try:
                part_no = str(row[mapped_columns['part_no']]).strip()
                if part_no.lower() in ('nan', '', 'none'):
                    continue
                item = {
                    'Part_No': part_no,
                    'Current_QTY': self.safe_float_convert(row[mapped_columns['current_qty']]),
                    'Current Inventory - VALUE': self.safe_float_convert(row.get(mapped_columns.get('stock_value', ''), 0)),
                    'Description': str(row.get(mapped_columns.get('description', ''), '')).strip(),
                    'UOM': str(row.get(mapped_columns.get('uom', ''), '')).strip(),
                    'Location': str(row.get(mapped_columns.get('location', ''), '')).strip(),
                    'Vendor_Code': str(row.get(mapped_columns.get('vendor_code', ''), '')).strip(),
                    'Batch': str(row.get(mapped_columns.get('batch', ''), '')).strip()
                }
                standardized_data.append(item)
                
            except Exception as e:
                if self.debug:
                    st.write(f"⚠️ Error processing row {i+1}: {e}")
                continue
        if self.debug:
            st.write(f"✅ Total standardized records: {len(standardized_data)}")
        return standardized_data
    
    def validate_inventory_against_pfep(self, inventory_data):
        """Validate inventory data against PFEP master data with normalized keys and warnings."""
        pfep_data = self.persistence.load_data_from_session_state('persistent_pfep_data')
        if not pfep_data:
            return {'is_valid': False, 'issues': ['No PFEP data available'], 'warnings': []}
        # Normalize part numbers
        def normalize(pn): return str(pn).strip().upper()
        pfep_df = pd.DataFrame(pfep_data)
        inventory_df = pd.DataFrame(inventory_data)
        pfep_df['Part_No'] = pfep_df['Part_No'].apply(normalize)
        inventory_df['Part_No'] = inventory_df['Part_No'].apply(normalize)

        pfep_parts = set(pfep_df['Part_No'])
        inventory_parts = set(inventory_df['Part_No'])

        issues = []
        warnings = []

        missing_parts = pfep_parts - inventory_parts
        extra_parts = inventory_parts - pfep_parts
        
        # Check for parts with zero quantity
        zero_qty_parts = inventory_df[inventory_df['Current_QTY'] == 0]['Part_No'].tolist()
        if zero_qty_parts:
            warnings.append(f"Parts with zero quantity: {len(zero_qty_parts)} parts")
        is_valid = len(issues) == 0
        return {
            'is_valid': is_valid,
            'issues': issues,
            'warnings': warnings,
            'pfep_parts_count': len(pfep_parts),
            'inventory_parts_count': len(inventory_parts),
            'matching_parts_count': len(pfep_parts & inventory_parts),
            'missing_parts_count': len(missing_parts),
            'extra_parts_count': len(extra_parts),
            'missing_parts_list': list(missing_parts),
            'extra_parts_list': list(extra_parts),
            'zero_qty_parts_list': zero_qty_parts
        }
        
    def admin_data_management(self):
        """
        Admin-only PFEP data management interface.
        Handles persistent settings, file upload, disk persistence, and locking.
        """
        st.header("🔧 Admin Dashboard - PFEP Data Management")
        
        # --- 1. GLOBAL SETTINGS (Auto-saves to Disk) ---
        st.markdown("### ⚙️ Global Analysis Settings")
        config_col1, config_col2 = st.columns(2)

        with config_col1:
            st.subheader("📐 Analysis Tolerance")
            
            # Load current value from session state (synced by initialize_session_state)
            current_tol = st.session_state.get('admin_tolerance', 30)
            options = [0, 10, 20, 30, 40, 50]
            
            # Find index of current value to keep dropdown synced
            try:
                idx = options.index(current_tol)
            except ValueError:
                idx = 3 # Fallback to 30%

            new_tolerance = st.selectbox(
                "Tolerance Zone (+/-)",
                options=options,
                index=idx,
                format_func=lambda x: f"{x}%",
                key="tol_selector_widget"
            )
            
            # If changed, update session AND disk immediately
            if new_tolerance != st.session_state.admin_tolerance:
                st.session_state.admin_tolerance = new_tolerance
                # FIX: Added st.session_state.company_id as the first argument
                DataPersistence.save_settings(
                    st.session_state.company_id, 
                    st.session_state.admin_tolerance, 
                    st.session_state.user_preferences['ideal_inventory_days']
                )
                st.success(f"✅ Tolerance locked at {new_tolerance}%")

        with config_col2:
            st.subheader("📅 Ideal Inventory Target")
            
            current_ideal = st.session_state.user_preferences.get('ideal_inventory_days', 30)
            
            new_ideal_days = st.number_input(
                "Ideal Inventory Days",
                min_value=1, max_value=365,
                value=int(current_ideal),
                step=1,
                key="ideal_days_input_widget"
            )
            
            # If changed, update session AND disk immediately
            if new_ideal_days != current_ideal:
                st.session_state.user_preferences['ideal_inventory_days'] = new_ideal_days
                # FIX: Added st.session_state.company_id as the first argument
                DataPersistence.save_settings(
                    st.session_state.company_id, 
                    st.session_state.admin_tolerance, 
                    st.session_state.user_preferences['ideal_inventory_days']
                )
                st.success(f"✅ Ideal Days locked at {new_ideal_days}")

        st.markdown("---")
        
        # --- 2. PFEP DATA STATUS & LOCKING ---
        pfep_locked = st.session_state.get('persistent_pfep_locked', False)
        
        if pfep_locked:
            st.warning("🔒 PFEP Master Data is currently LOCKED and active for Users.")
            col_a, col_b = st.columns([3, 1])
            with col_a:
                st.info("The Master Data is stored on the server disk. Users can now upload Inventory Dumps for analysis.")
            with col_b:
                if st.button("🔓 Unlock & Delete", type="secondary", help="Removes Master PFEP from disk and session"):
                    # Remove from disk
                    # Get the current company ID
                    comp_id = st.session_state.get('company_id')
                    # Remove from disk using the Class Name and passing the comp_id
                    DataPersistence.delete_from_disk(comp_id) 
                    # Clear session
                    st.session_state.persistent_pfep_data = None
                    st.session_state.persistent_pfep_locked = False
                    st.session_state.persistent_analysis_results = None
                    st.success("✅ PFEP deleted from server. Ready for new upload.")
                    st.rerun()
            
            # Show statistics of the locked data
            pfep_data = st.session_state.get('persistent_pfep_data')
            if pfep_data:
                self.display_pfep_data_preview(pfep_data)
            return

        # --- 3. UPLOAD INTERFACE (Visible only when UNLOCKED) ---
        st.subheader("📊 PFEP Master Data Upload")
        tab1, tab2, tab3 = st.tabs(["📁 Upload File", "🧪 Load Sample", "📋 Status"])
    
        with tab1:
            st.markdown("**Upload PFEP Excel/CSV File**")
            uploaded_file = st.file_uploader("Choose PFEP file", type=['xlsx', 'xls', 'csv'], key="admin_pfep_file_uploader")
            
            if uploaded_file is not None:
                try:
                    df = pd.read_csv(uploaded_file) if uploaded_file.name.endswith('.csv') else pd.read_excel(uploaded_file)
                    
                    # Standardize (row-by-row debugging removed as requested)
                    processed_data = self.standardize_pfep_data(df)
                    
                    if processed_data:
                        # Store in temp session state to prevent button-rerun errors
                        st.session_state['temp_standardized_pfep'] = processed_data
                        
                        st.success(f"✅ Processed {len(processed_data)} valid records.")
                        with st.expander("🔍 Preview processed data"):
                            st.dataframe(pd.DataFrame(processed_data).head(10))
                        
                        # THE SAVE & LOCK BUTTON
                        if st.button("💾 Save & Lock PFEP Permanently", type="primary"):
                            final_data = st.session_state.get('temp_standardized_pfep')
                            
                            # FIX: Get the current logged in company ID
                            comp_id = st.session_state.get('company_id')
                            
                            if final_data and comp_id:
                                current_now = datetime.now()
                                
                                # 1. Save to session state container
                                st.session_state.persistent_pfep_data = {
                                    'data': final_data,
                                    'timestamp': current_now
                                }
                                st.session_state.persistent_pfep_locked = True
                                
                                # 2. FIX: Pass comp_id as the first argument
                                DataPersistence.save_to_disk(comp_id, final_data, locked=True)
                                
                                st.success(f"✅ Master PFEP locked for {comp_id}. Date: {current_now.strftime('%d-%m-%Y')}")
                                
                                # Cleanup
                                del st.session_state['temp_standardized_pfep']
                                st.rerun()
                            else:
                                st.error("❌ Error: Missing data or Company ID. Please log in again.")
                except Exception as e:
                    st.error(f"❌ File processing error: {str(e)}")

        with tab2:
            st.info("🧪 Load internal samples for testing purposes.")
            if st.button("Load Demo PFEP"):
                sample = self.load_sample_pfep_data()
                st.session_state['temp_standardized_pfep'] = sample
                st.success("Demo data ready. Click 'Save' in Tab 1 to lock it.")

        with tab3:
            st.markdown("**Storage Status**")
            if not pfep_locked:
                st.info("No Master Data is currently locked. The system is waiting for Admin to upload the PFEP file.")
            else:
                st.success("Master Data is correctly synchronized.")

    def display_pfep_data_preview(self, pfep_data):
        """Displays statistics for the currently loaded PFEP data"""
        df_preview = pd.DataFrame(pfep_data)
        
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            st.metric("Total Parts", len(df_preview))
        with col2:
            v_col = 'Vendor Name' if 'Vendor Name' in df_preview.columns else 'Vendor_Name'
            unique_v = df_preview[v_col].nunique() if v_col in df_preview.columns else 0
            st.metric("Unique Vendors", unique_v)
        with col3:
            total_val = (df_preview['unit_price'] * df_preview['RM_IN_QTY']).sum() if 'unit_price' in df_preview.columns else 0
            st.metric("Total Master Value", f"₹{total_val/1000000:.1f}M")
        with col4:
            st.metric("Status", "Locked")

        with st.expander("📋 View Master Data"):
            st.dataframe(df_preview, use_container_width=True)
    
    def user_inventory_upload(self):
        st.header("📦 User: Inventory Upload & Analysis")
        
        pfep_container = st.session_state.get('persistent_pfep_data')
        pfep_locked = st.session_state.get('persistent_pfep_locked', False)
        
        if not pfep_locked or pfep_container is None:
            st.error("❌ PFEP Master Data is not locked by Admin.")
            return

        if isinstance(pfep_container, dict) and 'data' in pfep_container:
            pfep_data = pfep_container['data']
        else:
            pfep_data = pfep_container 

        uploaded_inv = st.file_uploader("Upload Current Inventory Dump", type=['xlsx', 'csv'], key="inv_up_user")
        
        if uploaded_inv:
            df_inv = pd.read_csv(uploaded_inv) if uploaded_inv.name.endswith('.csv') else pd.read_excel(uploaded_inv)
            standardized_inv = self.standardize_current_inventory(df_inv)
            
            if standardized_inv:
                # --- ADD THIS LINE HERE TO FIX THE SIDEBAR STATUS ---
                self.persistence.save_data_to_session_state('persistent_inventory_data', standardized_inv)
                # ----------------------------------------------------

                pfep_keys = set(str(item['Part_No']).strip().upper() for item in pfep_data)
                inv_keys = set(str(item['Part_No']).strip().upper() for item in standardized_inv)
                matches = pfep_keys.intersection(inv_keys)
                
                st.info(f"🔍 Match Check: Found {len(matches)} matching parts.")

                if st.button("🚀 Run Analysis", type="primary"):
                    current_tolerance = st.session_state.get('admin_tolerance', 30)
                    results = self.analyzer.analyze_inventory(pfep_data, standardized_inv, tolerance=current_tolerance)
                    if results:
                        st.session_state['persistent_analysis_results'] = results
                        st.success(f"✅ Analysis generated for {len(results)} items.")
                        st.rerun()

        if st.session_state.get('persistent_analysis_results'):
            st.markdown("---")
            self.display_analysis_results()
                
    def generate_ppt_report(self, analysis_results):
        """
        Generates a professional PowerPoint report.
        Slide 1: Cover with corrected alignment and branding.
        Slide 2: Performance Overview (matches target design).
        Slide 3: Status Breakdown with % labels and fixed logo overlap.
        """
        import io
        df = pd.DataFrame(analysis_results)
    
        # --- 1. Colors & Constants ---
        COLOR_DARK_TEXT = RGBColor(89, 81, 75)   # Dark brownish-grey
        COLOR_BADGE_BG = RGBColor(235, 230, 220) # Light tan
        COLOR_BLACK = RGBColor(0, 0, 0)  
        COLOR_DARK_TITLE = RGBColor(89, 81, 75)
    
        # --- 2. Calculations ---
        biz_unit = st.session_state.get('ppt_biz_unit', 'BUS PLANT').upper()
        ideal_days = st.session_state.user_preferences.get('ideal_inventory_days', 30)
        tolerance = st.session_state.admin_tolerance
        inv_date = datetime.now().strftime('%B %d, %Y')
    
        pfep_ts = self.persistence.get_data_timestamp('persistent_pfep_data')
        pfep_ref = pfep_ts.strftime('%d-%m-%Y %H:%M') if pfep_ts else "N/A"

        df['avg_numeric'] = df['AVG CONSUMPTION/DAY'].apply(self.safe_float_convert)
        df['unit_price_numeric'] = df['UNIT PRICE'].apply(self.safe_float_convert)
    
        total_qty = df['Current Inventory - Qty'].sum()
        total_avg_cons = df['avg_numeric'].sum()
        actual_inv_val = df['Current Inventory - VALUE'].sum()
    
        actual_inv_days = total_qty / total_avg_cons if total_avg_cons > 0 else 0
        actual_minr = actual_inv_val / 1_000_000
        ideal_val_total = (df['avg_numeric'] * ideal_days * df['unit_price_numeric']).sum()
        ideal_minr = ideal_val_total / 1_000_000
        variance_pct = ((actual_inv_days - ideal_days) / ideal_days * 100) if ideal_days > 0 else 0
    
        status_col = 'Status' if 'Status' in df.columns else 'INVENTORY REMARK STATUS'
        excess_minr = df[df[status_col] == 'Excess Inventory']['Stock Deviation Value'].sum() / 1_000_000
        short_minr = abs(df[df[status_col] == 'Short Inventory']['Stock Deviation Value'].sum()) / 1_000_000
        excess_days_val = max(0, actual_inv_days - (ideal_days * (1 + tolerance/100)))

        # --- 3. PPT Setup ---
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    
        logo_path = os.path.join(os.getcwd(), "Image.png")
        bg_path = os.path.join(os.getcwd(), "background.png")

        def add_logo_bottom_right(slide):
            if os.path.exists(logo_path):
                try:
                    slide.shapes.add_picture(logo_path, prs.slide_width - Inches(1.8), prs.slide_height - Inches(0.8), width=Inches(1.5))
                except: pass

        # ==========================================
        # SLIDE 1: COVER
        # ==========================================
        s1 = prs.slides.add_slide(prs.slide_layouts[6]) 
        if os.path.exists(bg_path):
            s1.shapes.add_picture(bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

        title_box = s1.shapes.add_textbox(0, Inches(3.0), prs.slide_width, Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "INVENTORY ANALYSER"
        p.font.bold = True; p.font.size = Pt(54); p.font.color.rgb = RGBColor(255, 255, 255); p.alignment = PP_ALIGN.CENTER

        tagline_box = s1.shapes.add_textbox(0, Inches(4.0), prs.slide_width, Inches(0.5))
        p2 = tagline_box.text_frame.paragraphs[0]
        p2.text = "Optimize Stock, Reduce Costs, and Forecast Smarter with AI"
        p2.font.size = Pt(22); p2.font.color.rgb = RGBColor(255, 255, 255); p2.alignment = PP_ALIGN.CENTER

        credit_box = s1.shapes.add_textbox(0, Inches(5.8), prs.slide_width, Inches(0.5))
        p3 = credit_box.text_frame.paragraphs[0]
        run1 = p3.add_run(); run1.text = "Developed by "
        run2 = p3.add_run(); run2.text = "Rimjhim Rani"; run2.font.bold = True
        run3 = p3.add_run(); run3.text = " | "
        run4 = p3.add_run(); run4.text = "Agilomatrix"; run4.font.bold = True
        # FIX: Changed CENTE to CENTER
        p3.font.size = Pt(20); p3.font.color.rgb = RGBColor(255, 255, 255); p3.alignment = PP_ALIGN.CENTER
        
        add_logo_bottom_right(s1)

        # ==========================================
        # SLIDE 2: PERFORMANCE OVERVIEW
        # ==========================================
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        add_logo_bottom_right(s2)

        # 1. THE BADGE (Top Left)
        badge = s2.shapes.add_shape(1, Inches(0.8), Inches(0.6), Inches(1.8), Inches(0.35))
        badge.fill.solid(); badge.fill.fore_color.rgb = COLOR_BADGE_BG; badge.line.width = 0
        btf = badge.text_frame
        btf.paragraphs[0].text = f"📊 {biz_unit}"
        btf.paragraphs[0].font.size = Pt(11); btf.paragraphs[0].font.color.rgb = COLOR_DARK_TEXT; btf.paragraphs[0].font.bold = True

        # 2. Main Title
        title_shape = s2.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(10), Inches(0.8))
        title_shape.text_frame.text = "Current Inventory Performance Overview"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36); title_shape.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_TEXT

        # 3. Description
        desc = s2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(1))
        dtf = desc.text_frame; dtf.word_wrap = True
        p = dtf.paragraphs[0]
        p.text = f"Snapshot analysis for the {biz_unit} facility as of {inv_date}, benchmarked against PFEP standards."
        p.font.size = Pt(14); p.font.color.rgb = COLOR_DARK_TEXT

        # 4. Big KPI Numbers
        kpis = [
            (f"{int(ideal_days)}", "Target Days", "Ideal inventory level"),
            (f"{actual_inv_days:.1f}", "Actual Days", "Current on-hand inventory"),
            (f"{variance_pct:,.0f}%", "Variance", "Over target inventory")
        ]
        for i, (val, lab, sub) in enumerate(kpis):
            box = s2.shapes.add_textbox(Inches(1 + i*4), Inches(3.0), Inches(3.5), Inches(1.5))
            tf = box.text_frame
            p1 = tf.paragraphs[0]; p1.text = val; p1.font.size = Pt(64); p1.alignment = PP_ALIGN.CENTER; p1.font.color.rgb = COLOR_DARK_TEXT
            p2 = tf.add_paragraph(); p2.text = lab; p2.font.size = Pt(20); p2.alignment = PP_ALIGN.CENTER; p2.font.color.rgb = COLOR_DARK_TEXT
            p3 = tf.add_paragraph(); p3.text = sub; p3.font.size = Pt(12); p3.alignment = PP_ALIGN.CENTER; p3.font.color.rgb = COLOR_DARK_TEXT

        # 5. Detail Sections (Bottom Bullets)
        l_box = s2.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(5), Inches(1.5))
        ltf = l_box.text_frame
        ltf.paragraphs[0].text = "Inventory Value Analysis"
        ltf.paragraphs[0].font.bold = True; ltf.paragraphs[0].font.size = Pt(18)
        
        for text in [f"Ideal Inventory: ₹{ideal_minr:,.2f} MINR", f"Actual Inventory: ₹{actual_minr:,.2f} MINR", f"Tolerance Level: {tolerance}%"]:
            p = ltf.add_paragraph(); p.text = text; p.font.size = Pt(14); p.level = 0

        r_box = s2.shapes.add_textbox(Inches(7.0), Inches(5.0), Inches(5), Inches(1.5))
        rtf = r_box.text_frame
        rtf.paragraphs[0].text = "Critical Deviations"
        rtf.paragraphs[0].font.bold = True; rtf.paragraphs[0].font.size = Pt(18)
        
        for text in [f"Excess: ₹{excess_minr:,.2f} MINR ({excess_days_val:.1f} days)", f"Shortage: ₹{short_minr:,.2f} MINR", "Action Required: Immediate rebalancing"]:
            p = rtf.add_paragraph(); p.text = text; p.font.size = Pt(14); p.level = 0

        # Footer
        footer = s2.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11), Inches(0.4))
        footer.text_frame.text = f"Analysis Reference: PFEP dated {pfep_ref} | Inventory captured {datetime.now().strftime('%d-%m-%Y')}"
        footer.text_frame.paragraphs[0].font.size = Pt(10); footer.text_frame.paragraphs[0].font.italic = True
        
        # ==========================================
        # SLIDE 3: STATUS BREAKDOWN (Percentage & Overlap Fixed)
        # ==========================================
        s3 = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 1. Title (Black)
        title_box3 = s3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6))
        p_t3 = title_box3.text_frame.paragraphs[0]
        p_t3.text = "Inventory Status Breakdown by Part Classification"
        p_t3.font.size = Pt(32); p_t3.font.color.rgb = COLOR_BLACK

        # 2. Insight Summary (Black)
        total_parts = len(df)
        short_count = len(df[df[status_col] == 'Short Inventory'])
        within_count = len(df[df[status_col] == 'Within Norms'])
        excess_count = len(df[df[status_col] == 'Excess Inventory'])
        
        short_pct = (short_count / total_parts) if total_parts > 0 else 0
        within_pct = (within_count / total_parts) if total_parts > 0 else 0
        excess_pct = (excess_count / total_parts) if total_parts > 0 else 0

        ins_box = s3.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.6))
        p_ins3 = ins_box.text_frame.paragraphs[0]
        p_ins3.text = f"Detailed analysis of {total_parts} total parts reveals critical supply chain imbalances. Over {short_pct:.0%} of parts are running below optimal levels."
        p_ins3.font.size = Pt(16); p_ins3.font.color.rgb = COLOR_BLACK

        # 3. Status Label (Centered above chart)
        st_lbl = s3.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(2.3), Inches(0.4))
        p_lbl = st_lbl.text_frame.paragraphs[0]
        p_lbl.text = "Status"; p_lbl.font.bold = True; p_lbl.font.size = Pt(24)
        p_lbl.font.color.rgb = COLOR_BLACK; p_lbl.alignment = PP_ALIGN.CENTER

        # 4. Pie Chart
        cd3 = ChartData()
        cd3.categories = ['Short', 'Within Norm', 'Excess']
        cd3.add_series('Inventory Status', (short_pct, within_pct, excess_pct))
        
        chart_shp = s3.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(4.5), Inches(2.2), Inches(4.3), Inches(3.2), cd3)
        chart = chart_shp.chart
        
        # FIX: Ensure chart doesn't generate its own title (prevents double Status text)
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(18); chart.legend.font.color.rgb = COLOR_BLACK
    
        points = chart.plots[0].series[0].points
        colors = [RGBColor(79, 129, 189), RGBColor(192, 80, 77), RGBColor(155, 187, 89)]
        for i, pt in enumerate(points):
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb = colors[i]
            
        # FIX: Force data labels to show ONLY Percentage as whole numbers (e.g. 58%)
        chart.plots[0].has_data_labels = True
        labels = chart.plots[0].data_labels
        labels.show_percentage = True
        labels.show_value = False       # This stops decimals like 0.58 from appearing
        labels.number_format = '0%'      # Ensures whole number formatting
        
        for pt in points:
            pt.data_label.font.size = Pt(14); pt.data_label.font.bold = True
            pt.data_label.font.color.rgb = COLOR_BLACK

        # 5. KPI Boxes (Black Font)
        within_val = df[df[status_col] == 'Within Norms']['Current Inventory - VALUE'].sum() / 1_000_000
        box_y = Inches(5.5); box_w = Inches(3.7); box_h = Inches(1.0)
    
        def add_kpi_box_black(slide, x, title, body):
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, box_y, box_w, box_h)
            shp.fill.background(); shp.line.color.rgb = RGBColor(180, 180, 180)
            tf = shp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p1 = tf.paragraphs[0]; p1.text = title; p1.font.bold = True; p1.font.size = Pt(18)
            p1.alignment = PP_ALIGN.CENTER; p1.font.color.rgb = COLOR_BLACK
            p2 = tf.add_paragraph(); p2.text = body; p2.font.size = Pt(10)
            p2.alignment = PP_ALIGN.CENTER; p2.font.color.rgb = COLOR_BLACK

        add_kpi_box_black(s3, Inches(0.8), "Within Tolerance", f"{within_count} parts ({within_pct:.0%}) within range (₹{within_val:.2f} MINR)")
        add_kpi_box_black(s3, Inches(4.8), "Excess Inventory", f"{excess_count} parts ({excess_pct:.0%}) excess by ₹{excess_minr:.2f} MINR")
        add_kpi_box_black(s3, Inches(8.8), "Critical Shortages", f"{short_count} parts ({short_pct:.0%}) short by ₹{short_minr:.2f} MINR")

        # 6. Recommendation Bar & Logo
        rec_box = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.75), Inches(10.5), Inches(0.40))
        rec_box.fill.solid(); rec_box.fill.fore_color.rgb = COLOR_BADGE_BG; rec_box.line.color.rgb = RGBColor(180, 180, 180)
        p_rec = rec_box.text_frame.paragraphs[0]
        p_rec.text = "Recommended Action: Prioritize procurement for shortage items while implementing excess stock reduction strategies."
        p_rec.font.size = Pt(11); p_rec.font.bold = True; p_rec.alignment = PP_ALIGN.CENTER; p_rec.font.color.rgb = COLOR_BLACK

        if os.path.exists(logo_path):
            try:
                s3.shapes.add_picture(logo_path, prs.slide_width - Inches(1.9), prs.slide_height - Inches(0.7), width=Inches(1.6))
            except: pass

        # ==========================================
        # SLIDE 4: DASHBOARD (Restored Title & No Targets)
        # ==========================================
        s4 = prs.slides.add_slide(prs.slide_layouts[6])
        
        # RESTORED TITLE BOX
        title_shp4 = s4.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(12), Inches(0.8))
        p_t4 = title_shp4.text_frame.paragraphs[0]
        p_t4.text = "Visual Inventory Analytics Dashboard"
        p_t4.font.size = Pt(40); p_t4.font.color.rgb = COLOR_BLACK

        # Descriptive subtitle
        sub_shp4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.6))
        sub_shp4.text_frame.paragraphs[0].text = f"Visual identification of inventory imbalances at the {biz_unit} facility."
        sub_shp4.text_frame.paragraphs[0].font.size = Pt(14); sub_shp4.text_frame.paragraphs[0].font.color.rgb = COLOR_BLACK

        # Data Prep (Pure Value Deviations in Lakhs)
        df_exc_top = df[df[status_col] == 'Excess Inventory'].sort_values('Stock Deviation Value', ascending=False).head(10)
        df_shrt_top = df[df[status_col] == 'Short Inventory'].sort_values('Stock Deviation Value', ascending=True).head(10)
        
        # Vendor Aggregation using standardized column 'Vendor Name'
        df_v_exc = df[df[status_col] == 'Excess Inventory'].groupby('Vendor Name')['Stock Deviation Value'].sum().reset_index().sort_values('Stock Deviation Value', ascending=False).head(10)
        df_v_shrt = df[df[status_col] == 'Short Inventory'].groupby('Vendor Name')['Stock Deviation Value'].sum().reset_index()
        df_v_shrt['Abs_Dev'] = abs(df_v_shrt['Stock Deviation Value'])
        df_v_shrt = df_v_shrt.sort_values('Abs_Dev', ascending=False).head(10)

        # Simple Bar Chart Helper (No Target comparison)
        def add_dashboard_bar(x_pos, title, categories, data, color):
            chart_data = ChartData()
            chart_data.categories = [str(c)[:10] for c in categories] # Short labels
            chart_data.add_series('Value (Lakhs)', [v / 100_000 for v in data])
            
            c_shp = s4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x_pos, Inches(2.2), Inches(3.1), Inches(3.2), chart_data)
            chart = c_shp.chart
            chart.has_title = True; chart.chart_title.text_frame.text = title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(9); chart.chart_title.text_frame.paragraphs[0].font.color.rgb = COLOR_BLACK
            chart.has_legend = False
            
            # Color formatting
            for pt in chart.series[0].points:
                pt.format.fill.solid(); pt.format.fill.fore_color.rgb = color
            chart.category_axis.tick_labels.font.size = Pt(7); chart.value_axis.tick_labels.font.size = Pt(7)

        # Draw the 4 Charts
        add_dashboard_bar(Inches(0.2), "Top 10 Excess Parts", df_exc_top['PART NO'], df_exc_top['Stock Deviation Value'], RGBColor(33, 150, 243))
        add_dashboard_bar(Inches(3.4), "Top 10 Short Parts", df_shrt_top['PART NO'], abs(df_shrt_top['Stock Deviation Value']), RGBColor(244, 67, 54))
        add_dashboard_bar(Inches(6.6), "Top 10 Excess Vendors", df_v_exc['Vendor Name'], df_v_exc['Stock Deviation Value'], RGBColor(33, 150, 243))
        add_dashboard_bar(Inches(9.8), "Top 10 Short Vendors", df_v_shrt['Vendor Name'], df_v_shrt['Abs_Dev'], RGBColor(244, 67, 54))

        # Bottom Recommendation Box
        bot_shp = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6))
        bot_shp.fill.solid(); bot_shp.fill.fore_color.rgb = COLOR_BADGE_BG
        p_bot = bot_shp.text_frame.paragraphs[0]
        p_bot.text = "📊 Visual analysis focused on magnitude of deviation to prioritize high-impact procurement actions."
        p_bot.font.size = Pt(14); p_bot.font.color.rgb = COLOR_BLACK; p_bot.alignment = PP_ALIGN.CENTER
        
        add_logo_bottom_right(s4)

        # ==========================================
        # SLIDE 5: CRITICAL ACTION ITEMS (MATCHING SCREENSHOT)
        # ==========================================
        s5 = prs.slides.add_slide(prs.slide_layouts[6])

        # 1. Main Title
        title_box5 = s5.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
        p_t5 = title_box5.text_frame.paragraphs[0]
        p_t5.text = "Critical Action Items: Top 10 Parts"
        p_t5.font.size = Pt(44); p_t5.font.color.rgb = COLOR_DARK_TITLE

        # 2. FIXED: Main Description with Word Wrap
        desc_box5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8))
        tf_d5 = desc_box5.text_frame
        tf_d5.word_wrap = True  # FIX: Ensures text wraps inside the box
        p_d5 = tf_d5.paragraphs[0]
        p_d5.text = "Focus areas for immediate inventory optimization. These parts represent the highest priority opportunities to reduce carrying costs and prevent production disruptions."
        p_d5.font.size = Pt(15); p_d5.font.color.rgb = COLOR_BLACK

        # 3. Sub-headers and Sub-descriptions
        # LEFT SIDE
        exc_title = s5.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.5), Inches(0.5))
        exc_title.text_frame.paragraphs[0].text = "Top 10 Excess Inventory Parts"
        exc_title.text_frame.paragraphs[0].font.size = Pt(24); exc_title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_TITLE
        
        exc_desc_box = s5.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(5.3), Inches(0.8))
        exc_desc_box.text_frame.word_wrap = True
        exc_desc_box.text_frame.paragraphs[0].text = "Overstocked items tying up working capital. Priority candidates for consumption acceleration."
        exc_desc_box.text_frame.paragraphs[0].font.size = Pt(12); exc_desc_box.text_frame.paragraphs[0].font.color.rgb = COLOR_BLACK

        # RIGHT SIDE
        shrt_title = s5.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(0.5))
        shrt_title.text_frame.paragraphs[0].text = "Top 10 Short Inventory Parts"
        shrt_title.text_frame.paragraphs[0].font.size = Pt(24); shrt_title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_TITLE
        
        shrt_desc_box = s5.shapes.add_textbox(Inches(7.2), Inches(2.7), Inches(5.3), Inches(0.8))
        shrt_desc_box.text_frame.word_wrap = True
        shrt_desc_box.text_frame.paragraphs[0].text = "Critical shortage items requiring expedited procurement action to maintain production continuity."
        shrt_desc_box.text_frame.paragraphs[0].font.size = Pt(12); shrt_desc_box.text_frame.paragraphs[0].font.color.rgb = COLOR_BLACK

        # 4. Table Helper with Qty column
        def create_action_table_with_qty(slide, x_pos, data_df, is_excess=True):
            # 6 rows, 5 columns (#, Part No, Description, Qty, Value)
            rows, cols = 6, 5
            table_shp = slide.shapes.add_table(rows, cols, x_pos, Inches(3.6), Inches(5.8), Inches(2.5))
            table = table_shp.table
            
            # Column Widths
            table.columns[0].width = Inches(0.4) # #
            table.columns[1].width = Inches(1.2) # Part No
            table.columns[2].width = Inches(2.3) # Description
            table.columns[3].width = Inches(0.9) # Qty
            table.columns[4].width = Inches(1.0) # Value Lakhs
            
            # Header
            headers = ["#", "Part No", "Description", "Qty", "Value(L)"]
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = h
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.text_frame.paragraphs[0].font.color.rgb = COLOR_BLACK
            
            # Fill Data (Top 5)
            for row_idx, (_, row_data) in enumerate(data_df.head(5).iterrows(), start=1):
                # Using Absolute Value for Qty display
                raw_qty = row_data.get('Stock Deviation Qty w.r.t Revised Norm', 0)
                disp_qty = abs(raw_qty)
                
                val = abs(row_data['Stock Deviation Value'])
                
                content = [
                    str(row_idx),
                    str(row_data['PART NO'])[:15],
                    str(row_data['PART DESCRIPTION'])[:22],
                    f"{disp_qty:,.0f}",
                    f"{val/100000:.2f}"
                ]
                for col_idx, text in enumerate(content):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = text
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(10)
                    p.font.color.rgb = COLOR_BLACK
                    # Align numbers center, description left
                    p.alignment = PP_ALIGN.CENTER if col_idx != 2 else PP_ALIGN.LEFT

        # Filter Data
        df_exc = df[df[status_col] == 'Excess Inventory'].sort_values('Stock Deviation Value', ascending=False)
        df_shrt = df[df[status_col] == 'Short Inventory'].sort_values('Stock Deviation Value', ascending=True)

        create_action_table_with_qty(s5, Inches(0.4), df_exc, is_excess=True)
        create_action_table_with_qty(s5, Inches(7.0), df_shrt, is_excess=False)
        add_logo_bottom_right(s5)
        
        ppt_out = io.BytesIO()
        prs.save(ppt_out)
        ppt_out.seek(0)
        return ppt_out
        
    def run(self):
        st.title("📊 Inventory Analyzer")
        st.markdown("<p style='font-size:18px; font-style:italic;'>Designed and Developed by Agilomatrix</p>", unsafe_allow_html=True)
        
        self.authenticate_user()

        # --- NEW LOGO UPLOADER SECTION ---
        st.sidebar.markdown("---")
        st.sidebar.subheader("🖼️ Branding")
        cust_logo = st.sidebar.file_uploader("Upload Customer Logo (Top Right)", type=['png', 'jpg', 'jpeg'])
        if cust_logo:
            st.session_state.customer_logo = cust_logo

        if st.session_state.user_role == "Admin":
            self.admin_data_management()
        elif st.session_state.user_role == "User":
            self.user_inventory_upload()

    def display_validation_results(self, validation_result):
        """Display inventory validation results"""
        st.markdown("**📋 Validation Results**")
        if validation_result['is_valid']:
            st.success("✅ Inventory data validation passed!")
        else:
            st.error("❌ Inventory data validation failed!")
            for issue in validation_result['issues']:
                st.error(f"• {issue}")
        
        # Display warnings
        if validation_result['warnings']:
            st.warning("⚠️ Validation Warnings:")
            for warning in validation_result['warnings']:
                st.warning(f"• {warning}")
        
        # Statistics
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("PFEP Parts", validation_result['pfep_parts_count'])
        with col2:
            st.metric("Inventory Parts", validation_result['inventory_parts_count'])
        with col3:
            st.metric("Matching Parts", validation_result['matching_parts_count'])
        
        # Detailed breakdown
        with st.expander("📊 Detailed Breakdown"):
            col1, col2 = st.columns(2)
            with col1:
                if validation_result['missing_parts_count'] > 0:
                    st.warning(f"Missing Parts ({validation_result['missing_parts_count']}):")
                    st.text("\n".join(validation_result['missing_parts_list'][:10]))
                    if len(validation_result['missing_parts_list']) > 10:
                        st.text(f"... and {len(validation_result['missing_parts_list']) - 10} more")
            
            with col2:
                if validation_result['extra_parts_count'] > 0:
                    st.info(f"Extra Parts ({validation_result['extra_parts_count']}):")
                    st.text("\n".join(validation_result['extra_parts_list'][:10]))
                    if len(validation_result['extra_parts_list']) > 10:
                        st.text(f"... and {len(validation_result['extra_parts_list']) - 10} more")
    
    def display_analysis_interface(self):
        """Main analysis interface for users"""
        st.subheader("📈 Inventory Analysis Results")
        # Get PFEP and Inventory data
        try:
            pfep_data = self.persistence.load_data_from_session_state('persistent_pfep_data')
            inventory_data = self.persistence.load_data_from_session_state('persistent_inventory_data')
        except Exception as e:
            st.error("❌ Error loading PFEP or Inventory data.")
            st.code(str(e))
            return
        if not pfep_data or not inventory_data:
            st.error("❌ Required data not available. Please upload PFEP and Inventory data first.")
            return
        # Get tolerance from admin settings
        tolerance = st.session_state.get('admin_tolerance', 30)
        st.info(f"📐 Analysis Tolerance: ±{tolerance}% (Set by Admin)")

        # Check if analysis needs to be performed or updated
        analysis_data = self.persistence.load_data_from_session_state('persistent_analysis_results')
        last_tolerance = st.session_state.get('last_analysis_tolerance', None)

        # Auto re-analyze if tolerance changed or no data exists
        if not analysis_data or last_tolerance != tolerance:
            st.info(f"🔄 Re-analyzing with ±{tolerance}% tolerance...")
            with st.spinner("Analyzing inventory..."):
                try:
                    analysis_results = self.analyzer.analyze_inventory(
                        pfep_data,
                        inventory_data,
                        tolerance=tolerance
                    )
                except Exception as e:
                    st.error("❌ Error during inventory analysis")
                    st.code(str(e))
                    return
            if analysis_results:
                self.persistence.save_data_to_session_state('persistent_analysis_results', analysis_results)
                st.session_state.last_analysis_tolerance = tolerance
                st.success("✅ Analysis completed successfully!")
                st.rerun()
            else:
                st.error("❌ Analysis failed. No results generated.")
                return
        # ✅ Use the full dashboard method
        try:
            self.display_analysis_results()
        except Exception as e:
            st.error("❌ Unexpected error during analysis results display")
            st.code(str(e))
            return
            
    def display_comprehensive_analysis(self, analysis_results):
        """Display comprehensive analysis results with enhanced features"""
        st.success(f"✅ Analysis Complete: {len(analysis_results)} parts analyzed")
        try:
            self.display_enhanced_summary_metrics(analysis_results)
        except Exception as e:
            st.error("❌ Error in Summary Metrics")
            st.code(str(e))
        try:
            self.display_enhanced_detailed_tables(analysis_results)
        except Exception as e:
            st.error("❌ Error in Detailed Tables")
            st.code(str(e))
        try:
            self.display_enhanced_analysis_charts(analysis_results)
        except Exception as e:
            st.error("❌ Error in Charts")
            st.code(str(e))
        try:
            self.display_enhanced_export_options(analysis_results)
        except Exception as e:
            st.error("❌ Error in Export Options")
            st.code(str(e))

    def display_enhanced_export_options(self, analysis_results):
        st.subheader("📤 Export Analysis Results")
        df = pd.DataFrame(analysis_results)
        col1, col2 = st.columns(2)
        with col1:
            excel_buffer = io.BytesIO()
            df.to_excel(excel_buffer, index=False)
            st.download_button("📊 Download Excel", data=excel_buffer.getvalue(), file_name="inventory.xlsx")
        with col2:
            if st.button("📑 Generate Professional PPT Report"):
                with st.spinner("Generating professional report..."):
                    # This calls the detailed function below
                    ppt_file = self.generate_ppt_report(analysis_results)
                    st.download_button(
                        label="📥 Download PPT Report",
                        data=ppt_file,
                        file_name=f"Inventory_Analysis_{datetime.now().strftime('%Y%m%d')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
        
    def display_enhanced_summary_metrics(self, analysis_results):
        """Enhanced summary metrics dashboard with Missing Parts Math"""
        st.header("📊 Executive Summary Dashboard")
        
        # 1. Fetch data for the Math
        pfep_container = st.session_state.get('persistent_pfep_data')
        pfep_list = pfep_container['data'] if isinstance(pfep_container, dict) else pfep_container
        inv_list = self.persistence.load_data_from_session_state('persistent_inventory_data')
        
        pfep_count = len(pfep_list) if pfep_list else 0
        inv_count = len(inv_list) if inv_list else 0
        missing_count = max(0, pfep_count - inv_count)

        # CSS Styling
        st.markdown("""
        <style>
        .metric-card {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 1.5rem; border-radius: 12px; margin: 0.5rem 0;
            box-shadow: 0 6px 15px rgba(0, 0, 0, 0.3); min-height: 140px;
        }
        .status-normal { background: linear-gradient(135deg, #4CAF50, #45a049); }
        .status-excess { background: linear-gradient(135deg, #2196F3, #1976D2); }
        .status-short { background: linear-gradient(135deg, #F44336, #D32F2F); }
        .status-total { background: linear-gradient(135deg, #FF9800, #F57C00); }
        .metric-value { color: white; font-weight: bold; font-size: 1.6rem; }
        .metric-label { color: #f0f0f0; font-size: 1.1rem; }
        .highlight-box {
            background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
            padding: 1.2rem; border-radius: 12px; color: white; margin: 1rem 0;
        }
        </style>
        """, unsafe_allow_html=True)

        df = pd.DataFrame(analysis_results)
        value_col = 'Stock Deviation Value'
        status_col = 'Status' if 'Status' in df.columns else 'INVENTORY REMARK STATUS'
        
        short_value = df[df[status_col] == 'Short Inventory'][value_col].sum() if not df.empty else 0
        excess_value = df[df[status_col] == 'Excess Inventory'][value_col].sum() if not df.empty else 0
        
        inv_val_col = next((c for c in ['Current Inventory - VALUE', 'Stock_Value'] if c in df.columns), None)
        total_stock_value = df[inv_val_col].sum() if inv_val_col else 0

        # math display in the highlight box
        st.markdown('<div class="highlight-box">', unsafe_allow_html=True)
        st.markdown(f"""
        ### 🎯 Key Inventory KPIs
        - **Total Master PFEP Parts**: {pfep_count:,}
        - **Parts in Current Inventory**: {inv_count:,}
        - **Missing Parts calculation**: {pfep_count:,} (PFEP) - {inv_count:,} (Inv) = **{missing_count:,} Missing Parts**
        - **Total Stock Value**: ₹{total_stock_value:,.0f}
        - **Shortage Impact**: ₹{abs(short_value):,.0f} | **Excess Impact**: ₹{excess_value:,.0f}
        """)
        st.markdown('</div>', unsafe_allow_html=True)

        # Status cards
        cols = st.columns(4)
        status_map = [
            ('Within Norms', '🟢 Within Norms', 'status-normal'),
            ('Excess Inventory', '🔵 Excess Inventory', 'status-excess'),
            ('Short Inventory', '🔴 Short Inventory', 'status-short')
        ]
        
        for i, (stat, label, css) in enumerate(status_map):
            filt = df[df[status_col] == stat]
            val = filt[inv_val_col].sum() if inv_val_col in filt.columns else 0
            cols[i].markdown(f"""<div class="metric-card {css}"><div class="metric-label">{label}</div>
            <div class="metric-value">{len(filt)} parts</div><div style='color:white'>₹{val:,.0f}</div></div>""", unsafe_allow_html=True)
        
        cols[3].markdown(f"""<div class="metric-card status-total"><div class="metric-label">📊 Total Analyzed</div>
        <div class="metric-value">{len(df)} parts</div><div style='color:white'>₹{total_stock_value:,.0f}</div></div>""", unsafe_allow_html=True)
            
    def display_enhanced_vendor_summary(self, analysis_results):
        """Enhanced vendor summary with better analytics"""
        st.header("🏢 Vendor Performance Analysis")
        df = pd.DataFrame(analysis_results)
        # Define vendor column only if one exist
        vendor_col = None
        if 'Vendor' in df.columns:
            vendor_col = 'Vendor'
        elif 'Vendor Name' in df.columns:
            vendor_col = 'Vendor Name'
        elif 'VENDOR' in df.columns:
            vendor_col = 'VENDOR'
        if vendor_col is None:
            st.warning("Vendor information not available in analysis data.")
            return  # Early exit if no vendor column found
        value_col = 'Current Inventory - VALUE'
        vendor_summary = {}
        for vendor in df[vendor_col].dropna().unique():
            vendor_data = df[df[vendor_col] == vendor]
            vendor_summary[vendor] = {
                'total_parts': len(vendor_data),
                'total_value': vendor_data[value_col].sum() if value_col in vendor_data.columns else 0,
                'short_parts': len(vendor_data[vendor_data['Status'] == 'Short Inventory']),
                'excess_parts': len(vendor_data[vendor_data['Status'] == 'Excess Inventory']),
                'normal_parts': len(vendor_data[vendor_data['Status'] == 'Within Norms']),
                'short_value': vendor_data[vendor_data['Status'] == 'Short Inventory'][value_col].sum() if value_col in vendor_data.columns else 0,
                'excess_value': vendor_data[vendor_data['Status'] == 'Excess Inventory'][value_col].sum() if value_col in vendor_data.columns else 0,
            }
            # Create vendor DataFrame OUTSIDE the loop
            vendor_df = pd.DataFrame([
                {
                    'Vendor': vendor,
                    'Total Parts': data['total_parts'],
                    'Short Inventory': data['short_parts'],
                    'Excess Inventory': data['excess_parts'],
                    'Within Norms': data['normal_parts'],
                    'Total Value (₹)': f"₹{data['total_value']:,.0f}",
                    'Performance Score': round((data['normal_parts'] / data['total_parts']) * 100, 1)
                    if data['total_parts'] > 0 else 0
                }
                for vendor, data in vendor_summary.items()
            ])
            if vendor_df.empty:
                st.warning("No vendor data available for analysis.")
                return
            def color_performance(val):
                if isinstance(val, (int, float)):
                    if val >= 80:
                        return 'background-color: #4CAF50; color: white'
                    elif val >= 60:
                        return 'background-color: #FF9800; color: white'
                    else:
                        return 'background-color: #F44336; color: white'
                return ''
            st.dataframe(
                vendor_df.style.map(color_performance, subset=['Performance Score']),
                use_container_width=True,
                hide_index=True
            )
            fig = px.bar(
                vendor_df.head(10),
                x='Vendor',
                y='Performance Score',
                title="Top 10 Vendor Performance Scores",
                color='Performance Score',
                color_continuous_scale='RdYlGn'
            )
            st.plotly_chart(fig, use_container_width=True)

    def display_enhanced_detailed_tables(self, analysis_results):
        """Display enhanced detailed tables with proper formatting"""
        st.header("📊 Detailed Analysis Tables")
        df = pd.DataFrame(analysis_results)
        # Create tabs for different views
        tab1, tab2, tab3, tab4 = st.tabs(["🔍 All Items", "🔴 Short Inventory", "🔵 Excess Inventory", "🟢 Within Norms"])
        with tab1:
            st.subheader("All Inventory Items")
            # Add search functionality
            search_term = st.text_input("🔍 Search parts (Part No, Description, Vendor):", key="search_all")
            display_df = df.copy()
            if search_term:
                search_mask = (
                    df['PART NO'].astype(str).str.contains(search_term, case=False, na=False) |
                    df.get('PART DESCRIPTION', pd.Series(dtype=str)).fillna('').str.contains(search_term, case=False, na=False) |
                    df.get('VENDOR', pd.Series(dtype=str)).fillna('').str.contains(search_term, case=False, na=False)
                )
                display_df = df[search_mask]
            # Select key columns for display
            display_columns = self._get_key_display_columns(display_df)
            if display_columns:
                styled_df = display_df[display_columns].style.format(self._get_column_formatters())
                st.dataframe(
                    styled_df,
                    use_container_width=True,
                    height=400,
                    hide_index=True
                )
            else:
                st.dataframe(display_df, use_container_width=True, height=400)
            st.info(f"📊 Showing {len(display_df)} of {len(df)} total items")
        with tab2:
            st.subheader("🔴 Short Inventory Items")
            short_items = df[df['Status'] == 'Short Inventory']
            if not short_items.empty:
                # Sort by impact/value
                value_col = self._get_value_column(short_items)
                if value_col:
                    short_items = short_items.sort_values(value_col, ascending=False)
                st.error(f"⚠️ {len(short_items)} items are short on inventory")
                # Add urgency classification
                if value_col:
                    short_items['Urgency'] = pd.cut(
                        short_items[value_col], 
                        bins=[0, 10000, 50000, float('inf')], 
                        labels=['Low', 'Medium', 'High'],
                        include_lowest=True
                    )
                display_columns = self._get_key_display_columns(short_items)
                if display_columns:
                    styled_df = short_items[display_columns].style.format(self._get_column_formatters())
                    st.dataframe(
                        styled_df,
                        use_container_width=True,
                        height=400,
                        hide_index=True
                    )
                else:
                    st.dataframe(short_items, use_container_width=True, height=400)
            else:
                st.success("✅ No items are currently short on inventory!")
        with tab3:
            st.subheader("🔵 Excess Inventory Items")
            excess_items = df[df['Status'] == 'Excess Inventory']
            if not excess_items.empty:
                # Sort by value
                value_col = self._get_value_column(excess_items)
                if value_col:
                    excess_items = excess_items.sort_values(value_col, ascending=False)
                st.warning(f"📦 {len(excess_items)} items have excess inventory")
                # Calculate potential savings
                if value_col:
                    total_excess_value = excess_items[value_col].sum()
                    st.metric("Total Excess Value", f"₹{total_excess_value:,.0f}")
                display_columns = self._get_key_display_columns(excess_items)
                if display_columns:
                    styled_df = excess_items[display_columns].style.format(self._get_column_formatters())
                    st.dataframe(
                        styled_df,
                        use_container_width=True,
                        height=400,
                        hide_index=True
                    )
                else:
                    st.dataframe(excess_items, use_container_width=True, height=400)
            else:
                st.success("✅ No items have excess inventory!")
        with tab4:
            st.subheader("🟢 Items Within Norms")
            normal_items = df[df['Status'] == 'Within Norms']
            if not normal_items.empty:
                st.success(f"✅ {len(normal_items)} items are within normal inventory levels")
                display_columns = self._get_key_display_columns(normal_items)
                if display_columns:
                    styled_df = normal_items[display_columns].style.format(self._get_column_formatters())
                    st.dataframe(
                        styled_df,
                        use_container_width=True,
                        height=400,
                        hide_index=True
                    )
                else:
                    st.dataframe(normal_items, use_container_width=True, height=400)
            else:
                st.warning("⚠️ No items are currently within normal inventory levels!")
        # Return the display dataframe for further use if needed
        display_columns = self._get_key_display_columns(df)
        return df[display_columns] if display_columns else df
        
    def _get_value_column(self, df):
        """Helper method to identify the main value column"""
        value_columns = ['Stock_Value', 'Current Inventory - VALUE', 'Current Inventory-VALUE']
        for col in value_columns:
            if col in df.columns:
                return col
        return None
        
    def _get_key_display_columns(self, df):
        """Helper method to select key columns for display"""
        # Define priority columns to show
        priority_columns = [
             # Part identification
            "PART NO",
            "PART DESCRIPTION",
            "Vendor Name",
            "Vendor_Code",
            "AVG CONSUMPTION/DAY",
            "RM IN DAYS",
            "RM Norm - In Qty",
            "Revised Norm Qty",
            "Lower Bound Qty",                # ✅ Added
            "Upper Bound Qty", 
            "UNIT PRICE",
            "Current Inventory - Qty",
            "Current Inventory - VALUE",
            "SHORT/EXCESS INVENTORY",
            "Stock Deviation Value",
            "INVENTORY REMARK STATUS"
        ]
        # Select columns that exist in the dataframe
        available_columns = []
        for col in priority_columns:
            if col in df.columns and col not in available_columns:
                available_columns.append(col)
        # Add any remaining important columns not in priority list
        for col in df.columns:
            if col not in available_columns and len(available_columns) < 10:
                available_columns.append(col)
        return available_columns # Limit to 10 c
        
    def _get_column_formatters(self,df=None):
        """Get column formatters for styling dataframes"""
        formatters = {}
        if df is not None:
            for col in df.columns:
                if 'VALUE' in col.upper() or 'PRICE' in col.upper() or 'COST' in col.upper():
                    formatters[col] = lambda x: f"₹{x:,.0f}" if pd.notnull(x) and isinstance(x, (int, float)) else str(x)
                elif 'QTY' in col.upper() or 'QUANTITY' in col.upper():
                    formatters[col] = lambda x: f"{x:,.0f}" if pd.notnull(x) and isinstance(x, (int, float)) else str(x)
                elif 'PERCENTAGE' in col.upper() or 'SCORE' in col.upper():
                    formatters[col] = lambda x: f"{x:.1f}%" if pd.notnull(x) and isinstance(x, (int, float)) else str(x)
        else:
            # Default formatters when no dataframe is provided
            formatters = {
                'VALUE(Unit Price* Short/Excess Inventory)': lambda x: f"₹{x:,.0f}" if pd.notnull(x) and isinstance(x, (int, float)) else str(x),
                'Current Inventory - VALUE': lambda x: f"₹{x:,.0f}" if pd.notnull(x) and isinstance(x, (int, float)) else str(x),
                'Current Inventory - Qty': lambda x: f"{x:,.0f}" if pd.notnull(x) and isinstance(x, (int, float)) else str(x),
                'UNIT PRICE': lambda x: f"₹{x:,.2f}" if pd.notnull(x) and isinstance(x, (int, float)) else str(x),
            }
        return formatters
    
    def display_overview_metrics(self, analysis_results):
        """Display key overview metrics"""
        st.header("📊 Inventory Overview")
        df = pd.DataFrame(analysis_results)
        # Create metrics columns
        col1, col2, col3, col4 = st.columns(4)
        with col1:
            total_parts = len(df)
            st.metric("Total Parts", f"{total_parts:,}")
        with col2:
            value_col = self._get_value_column(df)
            if value_col:
                total_value = df[value_col].sum()
                st.metric("Total Value", f"₹{total_value:,.0f}")
            else:
                st.metric("Total Value", "N/A")
        with col3:
            if 'Status' in df.columns:
                within_norms = (df['Status'] == 'Within Norms').sum()
                efficiency = (within_norms / total_parts * 100) if total_parts > 0 else 0
                st.metric("Efficiency", f"{efficiency:.1f}%", delta=f"{within_norms} parts")
            else:
                st.metric("Efficiency", "N/A")
        with col4:
            if 'Status' in df.columns:
                issues = len(df[df['Status'] != 'Within Norms'])
                st.metric("Issues", f"{issues:,}", delta="Needs attention" if issues > 0 else "All good")
            else:
                st.metric("Issues", "N/A")
                
    def display_top_parts_analysis(self, analysis_results):
        """Display top parts analysis by different criteria"""
        st.subheader("🏆 Top Parts Analysis")
        df = pd.DataFrame(analysis_results)
        value_col = self._get_value_column(df)
        if not value_col:
            st.warning("⚠️ No value column found for top parts analysis.")
            return
        col1, col2 = st.columns(2)
        with col1:
            st.subheader("💰 Highest Value Parts")
            top_value = df.nlargest(10, value_col)
            if not top_value.empty:
                display_cols = ['PART NO', 'PART DESCRIPTION', value_col, 'Status']
                available_cols = [col for col in display_cols if col in top_value.columns]
                st.dataframe(
                    top_value[available_cols].style.format(self._get_column_formatters(top_value)),
                    use_container_width=True
                )
        with col2:
            st.subheader("⚠️ Most Critical Issues")
            if 'Status' in df.columns:
                critical_issues = df[df['Status'] != 'Within Norms'].nlargest(10, value_col)
                if not critical_issues.empty:
                    display_cols = ['PART NO', 'PART DESCRIPTION', value_col, 'Status']
                    available_cols = [col for col in display_cols if col in critical_issues.columns]
                    st.dataframe(
                        critical_issues[available_cols].style.format(self._get_column_formatters(critical_issues)),
                        use_container_width=True
                    )
                else:
                    st.success("✅ No critical issues found!")
            else:
                st.info("ℹ️ Status information not available.")

    def create_enhanced_top_parts_chart(self, processed_data, status_filter, color, key, top_n=10):
        """Enhanced top parts chart with better visualization"""
        filtered_data = [
            item for item in processed_data 
            if item.get('Status') == status_filter or item.get('INVENTORY REMARK STATUS') == status_filter
        ]
        if not filtered_data:
            st.info(f"No {status_filter} parts found.")
            return

        top_parts = sorted(
            filtered_data,
            key=lambda x: x.get('Current Inventory - VALUE', 0),
            reverse=True
        )[:top_n]

        labels = [f"{item['PART NO']}<br>{item.get('PART DESCRIPTION', '')[:30]}..." for item in top_parts]
        values = [item.get('Current Inventory - VALUE', 0) for item in top_parts]
        variance_values = [item.get('VALUE(Unit Price* Short/Excess Inventory)', 0) for item in top_parts]

        fig = go.Figure()
        fig.add_trace(go.Bar(
            name='Stock Value',
            x=labels,
            y=values,
            marker_color=color,
            text=[f"₹{v:,.0f}" for v in values],
            textposition='auto',
        ))
        st.plotly_chart(fig, use_container_width=True)
        
    def apply_advanced_filters(self, df):
        """Apply advanced filters to dataframe"""
        filtered_df = df.copy()
        # Apply value filted
        if hasattr(st.session_state, 'value_filter') and 'Current Inventory - VALUE' in df.columns:
            min_val, max_val = st.session_state.value_filter
            filtered_df = filtered_df[
                (filtered_df['Current Inventory - VALUE'] >= min_val) & 
                (filtered_df['Current Inventory - VALUE'] <= max_val)
            ]
        # Apply quantity filter
        if hasattr(st.session_state, 'qty_filter') and 'Current Inventory - Qty' in df.columns:
            min_qty, max_qty = st.session_state.qty_filter
            filtered_df = filtered_df[
                (filtered_df['Current Inventory - Qty'] >= min_qty) & 
                (filtered_df['Current Inventory - Qty'] <= max_qty)
            ]
        # Apply vendor filter (FIXED: properly handle vendor column)
        if hasattr(st.session_state, 'vendor_filter'):
            vendor_col = None
            if 'Vendor' in df.columns:
                vendor_col = 'Vendor'
            elif 'Vendor Name' in df.columns:
                vendor_col = 'Vendor Name'
            if vendor_col and vendor_col in df.columns:
                filtered_df = filtered_df[filtered_df[vendor_col].isin(st.session_state.vendor_filter)]
        return filtered_df
            
    def display_trend_analysis(self, analysis_results):
        """
        UPDATED: Stock Coverage Report with Dynamic Category Selection.
        Defaults to 'Stock Less Than 3 Days' as requested.
        """
        st.header("📈 Stock Coverage Report")
        df = pd.DataFrame(analysis_results)
        
        # 1. Define the Categories and Binning Logic
        categories_order = [
            "Stock Out Items",
            "Stock Less Than 3 Days",
            "Stock 3-7 Days",
            "Stock 7-15 Days",
            "Stock 15-30 Days",
            "Stock 30-90 Days",
            "Stock 90-180 Days",
            "Stock 180 - 365 Days",
            "Stock More than Year",
            "No Consumption Parts"
        ]

        def get_coverage_category(row):
            qty = float(row.get('Current Inventory - Qty', 0) or 0)
            cons = float(row.get('AVG CONSUMPTION/DAY', 0) or 0)
            if cons <= 0: return "No Consumption Parts"
            if qty <= 0: return "Stock Out Items"
            
            days = qty / cons
            if days < 3: return "Stock Less Than 3 Days"
            if 3 <= days < 7: return "Stock 3-7 Days"
            if 7 <= days < 15: return "Stock 7-15 Days"
            if 15 <= days < 30: return "Stock 15-30 Days"
            if 30 <= days < 90: return "Stock 30-90 Days"
            if 90 <= days < 180: return "Stock 90-180 Days"
            if 180 <= days <= 365: return "Stock 180 - 365 Days"
            return "Stock More than Year"

        df['Coverage Category'] = df.apply(get_coverage_category, axis=1)

        # 2. Visual Summary (Table and Multi-Color Chart)
        summary_counts = df['Coverage Category'].value_counts().reindex(categories_order, fill_value=0).reset_index()
        summary_counts.columns = ['Coverage Category', 'No. of Items']
        
        sum_col1, sum_col2 = st.columns([1, 2])
        with sum_col1:
            st.subheader("📋 Summary")
            st.dataframe(summary_counts, hide_index=True, use_container_width=True)

        with sum_col2:
            colors = ['#d62728', '#ff7f0e', '#fdb462', '#e7ba52', '#bcbd22', '#1f77b4', '#17becf', '#9467bd', '#8c564b', '#7f7f7f']
            fig = px.bar(summary_counts, x='Coverage Category', y='No. of Items', color='Coverage Category',
                         color_discrete_sequence=colors, text='No. of Items', height=350)
            fig.update_layout(showlegend=False, xaxis_tickangle=-45, margin=dict(t=20, b=20, l=20, r=20))
            st.plotly_chart(fig, use_container_width=True)

        st.markdown("---")

        # 3. Dynamic Part List Selector (Drag Box)
        st.subheader("🔍 Detailed Part List by Category")
        
        # We set 'Stock Less Than 3 Days' as the default selection (index 1)
        selected_cat = st.selectbox(
            "Select Coverage Category to view items:",
            options=categories_order,
            index=1, 
            help="Choose a category to filter the table below. Default is 'Stock Less Than 3 Days'."
        )

        # Filter and show the list
        filtered_list = df[df['Coverage Category'] == selected_cat].copy()

        if not filtered_list.empty:
            st.info(f"Showing **{len(filtered_list)}** items for category: **{selected_cat}**")
            
            # Prepare columns for the detail view
            display_cols = ['PART NO', 'PART DESCRIPTION', 'Vendor Name', 'AVG CONSUMPTION/DAY', 'Current Inventory - Qty']
            
            # Add a 'Days of Stock' column for clarity (except for No Consumption)
            if selected_cat != "No Consumption Parts":
                filtered_list['Coverage (Days)'] = filtered_list.apply(
                    lambda r: round(float(r['Current Inventory - Qty'])/float(r['AVG CONSUMPTION/DAY']), 1) if float(r['AVG CONSUMPTION/DAY']) > 0 else 0, axis=1
                )
                display_cols.append('Coverage (Days)')

            st.dataframe(filtered_list[display_cols], use_container_width=True, height=400)
            
            # Category-specific Export
            csv_data = filtered_list[display_cols].to_csv(index=False).encode('utf-8')
            st.download_button(
                f"📥 Download {selected_cat} List",
                data=csv_data,
                file_name=f"{selected_cat.lower().replace(' ', '_')}.csv",
                mime="text/csv"
            )
        else:
            st.success(f"✅ No items found in the '{selected_cat}' category.")
            
    def export_comprehensive_report(self, analysis_results):
        """Export comprehensive analysis report"""
        try:
            df = pd.DataFrame(analysis_results)
            # Create Excel writer
            from io import BytesIO
            output = BytesIO()
        
            with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
                # Main analysis sheet
                df.to_excel(writer, sheet_name='Full Analysis', index=False)
                # Summary sheet
                summary_data = {
                    'Status': df['Status'].value_counts().index.tolist(),
                    'Count': df['Status'].value_counts().values.tolist(),
                    'Total Value': [df[df['Status'] == status]['Current Inventory - VALUE'].sum() 
                                    for status in df['Status'].value_counts().index]
                }
                pd.DataFrame(summary_data).to_excel(writer, sheet_name='Summary', index=False)
                # Critical items sheet
                critical_items = df[df['Current Inventory - VALUE'] > 100000]
                critical_items.to_excel(writer, sheet_name='Critical Items', index=False)
                # Download button
                st.download_button(
                    label="📥 Download Comprehensive Report",
                    data=output.getvalue(),
                    file_name=f"inventory_analysis_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
                st.success("✅ Comprehensive report prepared for download!")
        except Exception as e:
            st.error(f"❌ Export failed: {str(e)}")
            
    def export_critical_items(self, analysis_results):
        """Export only critical items"""
        try:
            df = pd.DataFrame(analysis_results)
            # Filter critical items
            critical_items = df[
                (df['Status'] != 'Within Norms') & 
                (df['Current Inventory - VALUE'] > st.session_state.get('critical_threshold', 100000))
            ]
            if critical_items.empty:
                st.warning("No critical items found based on current criteria.")
                return
            # Create CSV
            csv = critical_items.to_csv(index=False)
            st.download_button(
                label="📥 Download Critical Items Report",
                data=csv,
                file_name=f"critical_items_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv"
            )
            st.success(f"✅ Critical items report prepared! ({len(critical_items)} items)")
        except Exception as e:
            st.error(f"❌ Export failed: {str(e)}")
            
    def export_executive_summary(self, analysis_results):
        """Export executive summary"""
        try:
            df = pd.DataFrame(analysis_results)
            # Create executive summary data
            summary = {
                'Metric': [
                    'Total Parts Analyzed',
                    'Total Inventory Value (₹)',
                    'Parts Within Norms',
                    'Excess Inventory Parts',
                    'Short Inventory Parts',
                    'Inventory Efficiency (%)',
                    'Excess Value (₹)',
                    'Shortage Impact (₹)'
                ],
                'Value': [
                    len(df),
                    f"₹{df['Current Inventory - VALUE'].sum():,.0f}",
                    (df['Status'] == 'Within Norms').sum(),
                    (df['Status'] == 'Excess Inventory').sum(),
                    (df['Status'] == 'Short Inventory').sum(),
                    f"{(df['Status'] == 'Within Norms').mean() * 100:.1f}%",
                    f"₹{df[df['Status'] == 'Excess Inventory']['Current Inventory - VALUE'].sum():,.0f}",
                    f"₹{abs(df[df['Status'] == 'Short Inventory']['Stock Deviation Value'].sum()):,.0f}"
                ]
            }
            summary_df = pd.DataFrame(summary)
            csv = summary_df.to_csv(index=False)
            st.download_button(
                label="📥 Download Executive Summary",
                data=csv,
                file_name=f"executive_summary_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
                mime="text/csv"
            )
            st.success("✅ Executive summary prepared for download!")
        except Exception as e:
            st.error(f"❌ Export failed: {str(e)}")
            
    def export_custom_format(self, analysis_results, format_type):
        """Export in custom format"""
        try:
            df = pd.DataFrame(analysis_results)
            if format_type == "Excel (.xlsx)":
                self.export_comprehensive_report(analysis_results)
            elif format_type == "CSV (.csv)":
                csv = df.to_csv(index=False)
                st.download_button(
                    label="📥 Download CSV",
                    data=csv,
                    file_name=f"inventory_analysis_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.csv",
                    mime="text/csv"
                )
            elif format_type == "PDF Report":
                st.info("📄 PDF export functionality requires additional setup. Using CSV format instead.")
                self.export_custom_format(analysis_results, "CSV (.csv)")
            elif format_type == "PowerPoint Summary":
                st.info("📊 PowerPoint export functionality requires additional setup. Using Excel format instead.")
                self.export_comprehensive_report(analysis_results)
            st.success(f"✅ Export completed in {format_type} format!")
        except Exception as e:
            st.error(f"❌ Export failed: {str(e)}")
            
    def display_help_and_documentation(self):
        """Display help and documentation"""
        st.header("❓ Help & Documentation")
        with st.expander("📖 Understanding Analysis Results"):
            st.markdown("""
            ### Status Categories:
            - **🟢 Within Norms**: Inventory levels are optimal
            - **🔵 Excess Inventory**: Stock levels exceed requirements
            - **🔴 Short Inventory**: Stock levels are below minimum requirements
            
            ### Key Metrics:
            - **Stock Value**: Total monetary value of current inventory
            - **Variance Impact**: Financial impact of excess/shortage
            - **Performance Score**: Percentage of parts within norms
            """)
        with st.expander("🔧 Advanced Features"):
            st.markdown("""
            ### Filtering Options:
            - Use sidebar filters to focus on specific value ranges
            - Filter by vendor performance
            - Set critical value thresholds
            ### Export Options:
            - **Full Analysis**: Complete detailed report
            - **Critical Items**: High-value problem items only
            - **Executive Summary**: Key metrics for management
            """)
        with st.expander("💡 Best Practices"):
            st.markdown("""
            ### Optimization Tips:
            1. **Regular Monitoring**: Run analysis weekly/monthly
            2. **Vendor Performance**: Track vendor consistency
            3. **Critical Thresholds**: Adjust based on business needs
            4. **Action Items**: Follow up on recommendations
            5. **Trend Analysis**: Monitor patterns over time
            """)
        
                
    def display_advanced_filtering_options(self, analysis_results):
        """Display advanced filtering options in sidebar"""
        st.sidebar.header("🔍 Advanced Filters")
        df = pd.DataFrame(analysis_results)
        # Value range filter with proper handling of edge cases
        if 'Current Inventory - VALUE' in df.columns and not df['Current Inventory - VALUE'].empty:
            min_value = float(df['Current Inventory - VALUE'].min())
            max_value = float(df['Current Inventory - VALUE'].max())
        # Handle case where min and max are the same
            if min_value == max_value:
                if min_value == 0:
                    # If all values are 0, create a reasonable range
                    min_value = 0.0
                    max_value = 100000.0
                    st.sidebar.info("ℹ️ All stock values are 0. Using default range for filtering.")
                else:
                    # If all values are the same non-zero value, create a small range around i
                    range_buffer = max_value * 0.1 if max_value > 0 else 1000
                    min_value = max_value - range_buffer
                    max_value = max_value + range_buffer
                    st.sidebar.info(f"ℹ️ All stock values are {df['Current Inventory - VALUE'].iloc[0]:,.0f}. Adjusted range for filtering.")
            # Ensure min_value is always less than max_value
            if min_value >= max_value:
                max_value = min_value + 1000
            value_range = st.sidebar.slider(
                "Stock Value Range (₹)",
                min_value=min_value,
                max_value=max_value,
                value=(min_value, max_value),
                format="₹%.0f"
            )
        else:
            st.sidebar.warning("⚠️ Stock Value column not found or empty. Skipping value filter.")
            value_range = (0, 100000)  # Default range
        # Status filter
        if 'Status' in df.columns:
            status_options = df['Status'].unique().tolist()
            selected_statuses = st.sidebar.multiselect(
                "Filter by Status",
                options=status_options,
                default=status_options
            )
        else:
            selected_statuses = []
            st.sidebar.warning("⚠️ Status column not found.")
        # Category filter (if available)
        category_col = None
        if 'Category' in df.columns:
            category_col = 'Category'
        elif 'PART CATEGORY' in df.columns:
            category_col = 'PART CATEGORY'
        selected_categories = None
        if category_col and category_col in df.columns:
            categories = df[category_col].dropna().unique().tolist()
            if categories:
                selected_categories = st.sidebar.multiselect(
                    f"Filter by {category_col}",
                    options=categories,
                    default=categories
                )
            else:
                st.sidebar.info(f"ℹ️ No valid {category_col} values found.")
        # Vendor filter (if availablee)
        vendor_col = None
        if 'Vendor' in df.columns:
            vendor_col = 'Vendor'
        elif 'Vendor Name' in df.columns:
            vendor_col = 'Vendor Name'
        elif 'VENDOR' in df.columns:
            vendor_col = 'VENDOR'
        selected_vendors = None
        if vendor_col and vendor_col in df.columns:
            vendors = df[vendor_col].dropna().unique().tolist()
            if vendors:
                selected_vendors = st.sidebar.multiselect(
                    f"Filter by {vendor_col}",
                    options=vendors,
                    default=vendors
                )
            else:
                st.sidebar.info(f"ℹ️ No valid {vendor_col} values found.")
        # Critical threshold setting
        critical_threshold = st.sidebar.number_input(
            "Critical Value Threshold (₹)",
            min_value=0,
            value=100000,
            step=10000,
            help="Parts above this value are considered critical"
        )
        # Store filter values in session state
        st.session_state.filter_value_range = value_range
        st.session_state.filter_statuses = selected_statuses
        st.session_state.filter_categories = selected_categories
        st.session_state.filter_vendors = selected_vendors
        st.session_state.critical_threshold = critical_threshold
        # Display current filter summary
        with st.sidebar.expander("📋 Current Filters Summary"):
            st.write(f"**Value Range:** ₹{value_range[0]:,.0f} - ₹{value_range[1]:,.0f}")
            st.write(f"**Statuses:** {len(selected_statuses) if selected_statuses else 0} selected")
            if selected_categories:
                st.write(f"**Categories:** {len(selected_categories)} selected")
            if selected_vendors:
                st.write(f"**Vendors:** {len(selected_vendors)} selected")
            st.write(f"**Critical Threshold:** ₹{critical_threshold:,.0f}")

    def apply_advanced_filters(self, df):
        """Apply advanced filters to the dataframe with improved error handling"""
        filtered_df = df.copy()
        try:
            # Apply value range filter
            if (hasattr(st.session_state, 'filter_value_range') and 
                'Current Inventory - VALUE' in filtered_df.columns):
                    min_val, max_val = st.session_state.filter_value_range
                    filtered_df = filtered_df[
                        (filtered_df['Current Inventory - VALUE'] >= min_val) & 
                        (filtered_df['Current Inventory - VALUE'] <= max_val)
                    ]
            # Apply status filter
            if (hasattr(st.session_state, 'filter_statuses') and 
                st.session_state.filter_statuses and 
                'Status' in filtered_df.columns):
                    filtered_df = filtered_df[filtered_df['Status'].isin(st.session_state.filter_statuses)]
            # Apply category filter
            if (hasattr(st.session_state, 'filter_categories') and 
                st.session_state.filter_categories):
                    category_col = None
                    if 'Category' in filtered_df.columns:
                        category_col = 'Category'
                    elif 'PART CATEGORY' in filtered_df.columns:
                        category_col = 'PART CATEGORY'
                    if category_col:
                        filtered_df = filtered_df[filtered_df[category_col].isin(st.session_state.filter_categories)]
            # Apply vendor filter
            if (hasattr(st.session_state, 'filter_vendors') and 
                st.session_state.filter_vendors):
                    vendor_col = None
                    if 'Vendor' in filtered_df.columns:
                        vendor_col = 'Vendor'
                    elif 'Vendor Name' in filtered_df.columns:
                        vendor_col = 'Vendor Name'
                    elif 'VENDOR' in filtered_df.columns:
                        vendor_col = 'VENDOR'
                    if vendor_col:
                        filtered_df = filtered_df[filtered_df[vendor_col].isin(st.session_state.filter_vendors)]
            # Show filtering results
            original_count = len(df)
            filtered_count = len(filtered_df)
            if original_count != filtered_count:
                st.sidebar.success(f"✅ Filtered: {filtered_count:,} of {original_count:,} items")
            return filtered_df
        except Exception as e:
            st.sidebar.error(f"❌ Filter error: {str(e)}")
            st.sidebar.info("ℹ️ Returning unfiltered data")
            return df

    def generate_analysis_summary(self, analysis_results):
        """Generate a comprehensive analysis summary"""
        df = pd.DataFrame(analysis_results)
        summary = {
            'total_parts': len(df),
            'total_value': df['Current Inventory - VALUE'].sum(),
            'within_norms': (df['Status'] == 'Within Norms').sum(),
            'excess_inventory': (df['Status'] == 'Excess Inventory').sum(),
            'short_inventory': (df['Status'] == 'Short Inventory').sum(),
            'efficiency_rate': (df['Status'] == 'Within Norms').mean() * 100,
            'excess_value': df[df['Status'] == 'Excess Inventory']['Current Inventory - VALUE'].sum(),
            'shortage_impact': abs(df[df['Status'] == 'Short Inventory']['VALUE(Unit Price* Short/Excess Inventory)'].sum()),
            'critical_items': len(df[df['Current Inventory - VALUE'] > st.session_state.get('critical_threshold', 100000)]),
            'avg_stock_value': df['Current Inventory - VALUE'].mean()
        }
        return summary
        
    def main(self):
        """Main execution method for the inventory analyzer"""
        st.set_page_config(
            page_title="Advanced Inventory Analysis Dashboard",
            page_icon="📊",
            layout="wide",
            initial_sidebar_state="expanded"
        )
        # Initialize session state
        if 'analysis_complete' not in st.session_state:
            st.session_state.analysis_complete = False
        # Main application logic
        if not st.session_state.analysis_complete:
            # Show file upload and analysis interface
            self.display_file_upload_interface()
        else:
            # Show analysis results
            self.display_analysis_results()
        
        # Option to analyze new file
        if st.sidebar.button("🔄 Analyze New File", type="secondary"):
            st.session_state.analysis_complete = False
            st.rerun()
            
    def display_analysis_results(self):
        """Main method to display all analysis results"""
        analysis_results = st.session_state.get('persistent_analysis_results')
        
        if not analysis_results:
            st.warning("⚠️ No analysis results available yet. Please upload inventory and click 'Run Analysis'.")
            return

        # Define the results to be used by the charts below
        # This bypasses the sidebar filtering logic entirely
        filtered_results = analysis_results
    
        # Display main dashboard metrics and tables
        self.display_comprehensive_analysis(filtered_results)
    
        # Display the coverage report and charts
        st.markdown("---")
        self.display_trend_analysis(filtered_results)
    
        # Display help section
        st.markdown("---")
        self.display_help_and_documentation()
        
    def display_enhanced_analysis_charts(self, analysis_results):
        """Display enhanced visual summaries with Ideal Inventory Line Overlay"""
        st.subheader("📊 Enhanced Inventory Charts")
        
        # Add Unit Toggle AND Top N Slider
        col1, col2, col3 = st.columns([1, 2, 3])
        with col1:
            chart_unit = st.selectbox("Select Currency Unit:", ["Lakhs", "Millions"], key="chart_unit_selector")
        with col2:
            top_n = st.slider("Number of items to show:", min_value=5, max_value=50, value=10, step=5, key="top_n_slider")
        
        # ✅ Get Admin Configured Ideal Days
        ideal_days = st.session_state.get('user_preferences', {}).get('ideal_inventory_days', 30)
        
        if chart_unit == "Millions":
            divisor = 1_000_000
            suffix = "M"
            unit_name = "Millions"
        else:
            divisor = 100_000
            suffix = "L"
            unit_name = "Lakhs"
            
        df = pd.DataFrame(analysis_results)
        if df.empty:
            st.warning("⚠️ No data available for charts.")
            return
            
        # Identify value column
        value_col = next((col for col in ['Current Inventory - VALUE', 'Stock_Value', 'Current Inventory-VALUE'] if col in df.columns), None)
                
        if value_col and 'PART NO' in df.columns and 'PART DESCRIPTION' in df.columns:
            st.markdown(f"### 🏆 Top {top_n} Parts by Value")
            chart_data = df[df[value_col] > 0].sort_values(by=value_col, ascending=False).head(top_n).copy()
            chart_data['Value_Converted'] = chart_data[value_col] / divisor
            
            def calculate_ideal_metrics(row):
                try:
                    avg_consumption = float(row.get('AVG CONSUMPTION/DAY', 0) or 0)
                    unit_price = float(row.get('UNIT PRICE', 0) or 0)
                    current_qty = float(row.get('Current Inventory - Qty', 0) or 0)
                    ideal_qty = avg_consumption * ideal_days
                    ideal_value_converted = (ideal_qty * unit_price) / divisor
                    deviation_pct = ((current_qty - ideal_qty) / ideal_qty * 100) if ideal_qty > 0 else 0
                    return ideal_qty, ideal_value_converted, deviation_pct
                except: return 0, 0, 0

            metrics = chart_data.apply(calculate_ideal_metrics, axis=1, result_type='expand')
            chart_data['Ideal_Qty'], chart_data['Ideal_Value_Converted'], chart_data['Deviation_Pct'] = metrics[0], metrics[1], metrics[2]
            chart_data['Part'] = chart_data.apply(lambda row: f"{row['PART DESCRIPTION']}\n({row['PART NO']})", axis=1)
            chart_data['Inventory_Status'] = chart_data['Status'] if 'Status' in chart_data.columns else 'Within Norms'
            
            color_map = {"Excess Inventory": "#2196F3", "Short Inventory": "#F44336", "Within Norms": "#4CAF50"}
            chart_data['HOVER_TEXT'] = chart_data.apply(lambda row: (
                f"Description: {row['PART DESCRIPTION']}<br>Part No: {row['PART NO']}<br><b>Current Status: {row['Inventory_Status']}</b><br>"
                f"--------------------------------<br>Current Qty: {row.get('Current Inventory - Qty', 0):,.0f}<br>Current Value: ₹{row[value_col]:,.0f}<br>"
                f"--------------------------------<br><b>Ideal Target ({ideal_days} Days):</b><br>Ideal Qty: {row['Ideal_Qty']:,.0f}<br><b>Deviation: {row['Deviation_Pct']:.1f}%</b>"
            ), axis=1)

            fig1 = go.Figure()
            for i, row in chart_data.iterrows():
                fig1.add_trace(go.Bar(x=[row['Part']], y=[row['Value_Converted']], name=row['Inventory_Status'], marker_color=color_map.get(row['Inventory_Status'], "#4CAF50"), customdata=[row['HOVER_TEXT']], hovertemplate='<b>%{x}</b><br>%{customdata}<extra></extra>', showlegend=False))
            
            for status, color in color_map.items():
                fig1.add_trace(go.Bar(x=[None], y=[None], name=status, marker_color=color, showlegend=True))
                
            fig1.add_trace(go.Scatter(x=chart_data['Part'], y=chart_data['Ideal_Value_Converted'], mode='lines+markers', name=f'Ideal Inventory', line=dict(color='black', width=1.5), marker=dict(symbol='circle', size=5, color='black'), hovertemplate=f"<b>Ideal Target</b><br>Value: %{{y:.2f}} {suffix}<extra></extra>"))
            fig1.update_layout(title=f"Top {top_n} Parts Actual Value vs Ideal Target", xaxis_title="Parts", yaxis_title=f"Stock Value (₹ {unit_name})", xaxis_tickangle=-45, yaxis=dict(tickformat=',.1f', ticksuffix=suffix), height=600, showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
            st.plotly_chart(fig1, use_container_width=True)

        # ✅ 2. Vendor vs Value
        vendor_col = next((col for col in ['Vendor', 'Vendor Name', 'VENDOR'] if col in df.columns), None)
        if vendor_col and value_col:
            st.markdown(f"### 🏢 Top {top_n} Vendors by Value")
            vendor_data = []
            for vendor_name, vendor_group in df[df[value_col] > 0].groupby(vendor_col):
                total_value = vendor_group[value_col].sum()
                vendor_ideal_value_sum = sum(float(v.get('AVG CONSUMPTION/DAY', 0) or 0) * ideal_days * float(v.get('UNIT PRICE', 0) or 0) for _, v in vendor_group.iterrows())
                v_status = vendor_group['Status'].value_counts().index[0] if 'Status' in vendor_group.columns and not vendor_group['Status'].empty else 'Within Norms'
                vendor_data.append({vendor_col: vendor_name, value_col: total_value, 'Ideal_Value': vendor_ideal_value_sum, 'Vendor_Status': v_status})
            
            vendor_df = pd.DataFrame(vendor_data).sort_values(by=value_col, ascending=False).head(top_n)
            if not vendor_df.empty:
                vendor_df['Value_Converted'] = vendor_df[value_col] / divisor
                vendor_df['Ideal_Value_Converted'] = vendor_df['Ideal_Value'] / divisor
                vendor_df['Deviation_Pct'] = vendor_df.apply(lambda r: ((r[value_col] - r['Ideal_Value']) / r['Ideal_Value'] * 100) if r['Ideal_Value'] > 0 else 0, axis=1)
                
                vendor_df['HOVER_TEXT'] = vendor_df.apply(lambda row: (f"Vendor: {row[vendor_col]}<br>Actual Value: ₹{row[value_col]:,.0f}<br>Ideal Value: ₹{row['Ideal_Value']:,.0f}<br>Deviation: {row['Deviation_Pct']:.1f}%<br>Status: {row['Vendor_Status']}"), axis=1)
                
                fig2 = go.Figure()
                for i, row in vendor_df.iterrows():
                    fig2.add_trace(go.Bar(x=[row[vendor_col]], y=[row['Value_Converted']], name=row['Vendor_Status'], marker_color=color_map.get(row['Vendor_Status'], "#4CAF50"), customdata=[row['HOVER_TEXT']], hovertemplate='%{customdata}<extra></extra>', showlegend=False))
                
                fig2.add_trace(go.Scatter(x=vendor_df[vendor_col], y=vendor_df['Ideal_Value_Converted'], mode='lines+markers', name=f'Ideal Inventory', line=dict(color='black', width=1.5), marker=dict(symbol='circle', size=5, color='black')))
                fig2.update_layout(title=f'Top {top_n} Vendors Actual Value vs Ideal Target', xaxis_title="Vendors", yaxis_title=f"Inventory Value (₹ {unit_name})", xaxis_tickangle=-45, yaxis=dict(tickformat=',.1f', ticksuffix=suffix), height=600, showlegend=True, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
                st.plotly_chart(fig2, use_container_width=True)

        # ✅ 3. Top N Parts by Inventory Status
        try:
            st.markdown(f"## 🧩 Top {top_n} Parts by Inventory Status") 
            status_colors = {"Excess Inventory": "#2196F3", "Short Inventory": "#F44336"}
            for status, label, color in [("Excess Inventory", f"🔵 Top {top_n} Excess Inventory Parts", status_colors["Excess Inventory"]), ("Short Inventory", f"🔴 Top {top_n} Short Inventory Parts", status_colors["Short Inventory"])]:
                st.subheader(label)
                st.markdown(f'<div class="graph-description">Top {top_n} parts with highest {status.lower()} value compared to target.</div>', unsafe_allow_html=True)
                
                status_df = df[df['INVENTORY REMARK STATUS'] == status].copy()
                if status == "Excess Inventory":
                    status_df = status_df[status_df['Stock Deviation Value'] > 0].sort_values(by='Stock Deviation Value', ascending=False).head(top_n)
                else:
                    status_df = status_df[status_df['Stock Deviation Value'] < 0]
                    status_df['Abs_Dev'] = abs(status_df['Stock Deviation Value'])
                    status_df = status_df.sort_values(by='Abs_Dev', ascending=False).head(top_n)
                
                if status_df.empty: continue

                # Logic Change: Bar shows Total Current Value
                status_df['Value_Converted'] = status_df[value_col] / divisor
                status_df['Ideal_Val_Converted'] = status_df.apply(lambda r: (float(r.get('AVG CONSUMPTION/DAY', 0) or 0) * ideal_days * float(r.get('UNIT PRICE', 0) or 0)) / divisor, axis=1)
                status_df['PART_DESC_NO'] = status_df['PART DESCRIPTION'].astype(str) + " (" + status_df['PART NO'].astype(str) + ")"
                
                status_df['HOVER_TEXT'] = status_df.apply(lambda row: (
                    f"Description: {row.get('PART DESCRIPTION', 'N/A')}<br>Part No: {row.get('PART NO')}<br>"
                    f"Actual Value: ₹{row[value_col]:,.0f}<br>Ideal Target: ₹{(row['Ideal_Val_Converted'] * divisor):,.0f}<br>"
                    f"{status.split()[0]} Amount: ₹{abs(row['Stock Deviation Value']):,.0f}"
                ), axis=1)
                
                fig3 = go.Figure()
                fig3.add_trace(go.Bar(x=status_df['PART_DESC_NO'], y=status_df['Value_Converted'], marker_color=color, name='Actual Value', customdata=status_df['HOVER_TEXT'], hovertemplate='<b>%{x}</b><br>%{customdata}<extra></extra>'))
                fig3.add_trace(go.Scatter(x=status_df['PART_DESC_NO'], y=status_df['Ideal_Val_Converted'], mode='lines+markers', name='Ideal Target', line=dict(color='black', width=1.5), marker=dict(symbol='circle', size=5, color='black')))
                fig3.update_layout(title=f"Top {top_n} {status} Parts: Actual vs Target", xaxis_tickangle=-45, yaxis_title=f"Value (₹ {unit_name})", yaxis=dict(tickformat=',.1f', ticksuffix=suffix), height=600, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
                st.plotly_chart(fig3, use_container_width=True, key=f"{status.lower().replace(' ', '_')}_parts")
        except Exception as e: st.error("❌ Error in Parts Status Charts")

        # ✅ 4. Top N Vendors by Inventory Status
        try:
            st.markdown(f"## 🏢 Top {top_n} Vendors by Inventory Status")
            for status, metric_name, color in [("Excess Inventory", "Excess Value Above Norm", "#2196F3"), ("Short Inventory", "Short Value Below Norm", "#F44336")]:
                status_df = df[df['INVENTORY REMARK STATUS'] == status].copy()
                if status_df.empty: continue

                vendor_stats = []
                for vendor, group in status_df.groupby(vendor_col):
                    actual_sum = group[value_col].sum()
                    ideal_sum = sum(float(r.get('AVG CONSUMPTION/DAY', 0) or 0) * ideal_days * float(r.get('UNIT PRICE', 0) or 0) for _, r in group.iterrows())
                    dev_sum = abs(group['Stock Deviation Value'].sum())
                    vendor_stats.append({'Vendor': vendor, 'Actual_Value': actual_sum, 'Ideal_Value': ideal_sum, 'Dev_Value': dev_sum})

                v_df = pd.DataFrame(vendor_stats).sort_values(by='Dev_Value', ascending=False).head(top_n)
                v_df['Actual_Converted'] = v_df['Actual_Value'] / divisor
                v_df['Ideal_Converted'] = v_df['Ideal_Value'] / divisor

                fig4 = go.Figure()
                fig4.add_trace(go.Bar(x=v_df['Vendor'], y=v_df['Actual_Converted'], name='Actual Value', marker_color=color, hovertemplate=f'<b>{{x}}</b><br>Actual: ₹{{y:,.1f}} {suffix}<extra></extra>'))
                fig4.add_trace(go.Scatter(x=v_df['Vendor'], y=v_df['Ideal_Converted'], mode='lines+markers', name='Ideal Target', line=dict(color='black', width=1.5), marker=dict(symbol='circle', size=6, color='black'), hovertemplate=f'<b>Ideal Target</b><br>Value: ₹{{y:,.1f}} {suffix}<extra></extra>'))
                fig4.update_layout(title=f"Top {top_n} Vendors: {status} Actual vs Target", xaxis_title="Vendor", yaxis_title=f"Value (₹ {unit_name})", xaxis_tickangle=-45, yaxis=dict(tickformat=',.1f', ticksuffix=suffix), height=600, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
                st.plotly_chart(fig4, use_container_width=True, key=f"vendor_{status}_chart")
                
        except Exception as e: 
            st.error("❌ Error in Vendor Status Charts")
            
if __name__ == "__main__":
    app = InventoryManagementSystem()
    app.run()  # This runs the full dashboard
